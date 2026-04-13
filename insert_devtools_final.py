"""
insert_devtools_final.py
1. Removes the 11 slides added previously (indices 30-40)
2. Inserts 2 Chairman/HOD-level dev slides after slide 9 (Agentic Review Pipeline)
   - "Developer Admin Dashboard"
   - "System Quality & Analytics Monitoring"
Design matches existing: #12141A bg, #5B8DFF header bar, KAVIN logo extracted.
"""

import io, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

PPT_PATH = Path(r"D:\KavinBase\KavinBase_Presentation_Updated.pptx")

# ── Exact palette from existing slides ─────────────────────────────────────────
BG      = RGBColor(0x12, 0x14, 0x1A)
HEADER  = RGBColor(0x5B, 0x8D, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SUB     = RGBColor(0xDF, 0xE9, 0xFF)
BODY    = RGBColor(0xCC, 0xD6, 0xF1)
DIM     = RGBColor(0x70, 0x80, 0xA0)

BLUE    = RGBColor(0x5B, 0x8D, 0xFF)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
GREEN   = RGBColor(0x4A, 0xDE, 0x80)
RED     = RGBColor(0xFF, 0x45, 0x45)
CYAN    = RGBColor(0x22, 0xD3, 0xEE)
ORANGE  = RGBColor(0xFB, 0xBF, 0x24)

FILL_BLUE   = RGBColor(0x06, 0x12, 0x28)
FILL_PURPLE = RGBColor(0x10, 0x08, 0x28)
FILL_GREEN  = RGBColor(0x05, 0x18, 0x10)
FILL_RED    = RGBColor(0x1E, 0x06, 0x06)
FILL_ORANGE = RGBColor(0x1E, 0x12, 0x04)
FILL_CYAN   = RGBColor(0x04, 0x14, 0x20)

FONT = "Calibri"
SW, SH = 13.33, 7.50
HDR_H  = 1.10

_LOGO: bytes = b""


# ── Shape helpers ──────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=FILL_BLUE, border=None, bpt=1.5):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def tb(slide, x, y, w, h):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.text_frame.word_wrap = True
    return t


def para(tf, text, size=13, color=BODY, bold=False, italic=False,
         align=PP_ALIGN.LEFT, spc_after=1):
    # pick first empty para or add new
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment    = align
    p.space_after  = Pt(spc_after)
    r = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.name = FONT; r.font.color.rgb = color
    return p


# ── Chrome: background + blue header + logo ────────────────────────────────────

def chrome(slide, title, subtitle):
    rect(slide, 0, 0, SW, SH, fill=BG)                # dark bg
    rect(slide, 0, 0, SW, HDR_H, fill=HEADER)          # blue header bar

    t1 = tb(slide, 0.35, 0.10, SW - 2.55, 0.62)
    t1.text_frame.margin_top = Pt(0)
    para(t1.text_frame, title, size=26, color=WHITE, bold=True)

    t2 = tb(slide, 0.35, 0.68, SW - 2.55, 0.36)
    t2.text_frame.margin_top = Pt(0)
    para(t2.text_frame, subtitle, size=13, color=SUB, italic=False)

    if _LOGO:
        slide.shapes.add_picture(
            io.BytesIO(_LOGO),
            Inches(SW - 2.10), Inches(0.20),
            width=Inches(1.83), height=Inches(0.54)
        )
    return HDR_H + 0.18   # content y-start


# ── Card ───────────────────────────────────────────────────────────────────────

def card(slide, x, y, w, h, accent, fill, title, lines, icon=""):
    rect(slide, x, y, w, h, fill=fill, border=accent, bpt=1.8)
    # Title
    title_tb = tb(slide, x + 0.10, y + 0.09, w - 0.20, 0.44)
    title_tb.text_frame.margin_top = Pt(0)
    lbl = f"{icon}  {title}" if icon else title
    para(title_tb.text_frame, lbl, size=14, color=accent, bold=True)
    # Separator
    rect(slide, x + 0.08, y + 0.48, w - 0.16, 0.013, fill=accent)
    # Body
    body_tb = tb(slide, x + 0.12, y + 0.52, w - 0.22, h - 0.58)
    body_tb.text_frame.word_wrap = True
    body_tb.text_frame.margin_top = Pt(2)
    for line in lines:
        sz  = 11 if line.startswith("  ") else 12
        col = DIM  if line.startswith("  ") else BODY
        para(body_tb.text_frame, line, size=sz, color=col, spc_after=2)


# ── Slide removal ──────────────────────────────────────────────────────────────

def remove_slides_from(prs, start_idx):
    """Remove all slides from start_idx (0-based) to end."""
    sldIdLst = prs.slides._sldIdLst
    r_id_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    all_sldId = list(sldIdLst)

    for sldId in all_sldId[start_idx:]:
        rId = sldId.get(r_id_attr)
        sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass

    print(f"  Removed slides {start_idx + 1} onwards. Remaining: {len(prs.slides)}")


# ── Slide position move ────────────────────────────────────────────────────────

def move_slide_to(prs, from_idx, to_idx):
    """Move slide at from_idx to to_idx (both 0-based)."""
    sldIdLst = prs.slides._sldIdLst
    all_elem  = list(sldIdLst)
    elem = all_elem[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_admin_dashboard(slide):
    """
    Slide A: Developer Admin Dashboard — 2x2 card grid.
    High-level, Chairman/HOD focused.
    """
    cy = chrome(
        slide,
        "Developer Admin Dashboard",
        "Enterprise-grade control panel for live system management  |  "
        "Admin-gated  |  No server restart required for any operation"
    )

    gw  = (SW - 0.58) / 2
    gh  = (SH - cy - 0.32) / 2
    gap = 0.18

    cards = [
        (BLUE,   FILL_BLUE,   "AI",    "AI Model Management",
         ["Switch between three AI model modes at runtime",
          "  base  — Full HuggingFace transformer (GPU)",
          "  lite   — GGUF quantised model  (CPU only)",
          "  net   — Cloud API (OpenAI / Gemini)",
          "",
          "Download, register and test new models via API",
          "Hot-swap the active model without any downtime",
          "Model readiness checked per mode with error detail"]),

        (PURPLE, FILL_PURPLE, "USR",   "User Administration",
         ["Create users with isolated storage resources",
          "  Each user gets a dedicated PostgreSQL database",
          "  Personal MinIO bucket for document storage",
          "  Scoped Redis namespace for session data",
          "",
          "Enable / disable user accounts instantly",
          "Role-based access control (admin / user)",
          "Admin password reset without old credentials"]),

        (GREEN,  FILL_GREEN,  "MON",   "System Analytics & Monitoring",
         ["Live dashboard: GET /devtools/mvp/overview",
          "  Configurable time window: 1 to 168 hours",
          "",
          "Ingestion funnel: processing / ready / error counts",
          "RAG quality score: high / medium / low distribution",
          "Response latency and semantic cache hit rate",
          "User feedback tracking: positive vs negative ratio",
          "Runtime health: GPU VRAM, queues, worker status"]),

        (ORANGE, FILL_ORANGE, "DBS",   "Data Visibility & Controls",
         ["Live browser across all storage layers:",
          "  RAG DB    — pgvector semantic index",
          "  Chat DB   — conversation history",
          "  Redis      — session state & cache keys",
          "  MinIO      — uploaded document objects",
          "",
          "Paginated record inspection (no raw SQL needed)",
          "Selective reset controls with safety confirmation",
          "Per-session & per-user RAG toggle overrides"]),
    ]

    for i, (ac, fi, ico, ttl, buls) in enumerate(cards):
        ci = i % 2; ri = i // 2
        card(slide,
             0.20 + ci * (gw + gap),
             cy  + ri * (gh + gap * 0.65),
             gw, gh, ac, fi, ttl, buls, ico)


def build_analytics(slide):
    """
    Slide B: System Quality & Analytics Monitoring — 4 metric cards.
    Shows the system is self-monitoring and enterprise-ready.
    """
    cy = chrome(
        slide,
        "System Quality & Analytics Monitoring",
        "Real-time operational metrics proving system reliability  |  "
        "GET /devtools/mvp/overview  |  Configurable 1-168 hour window"
    )

    cw  = (SW - 0.58) / 4
    ch  = SH - cy - 0.28
    gap = 0.18

    cards = [
        (RED,    FILL_RED,    "RT",  "Runtime Health",
         ["GPU health check",
          "  Device name + VRAM usage",
          "  Free vs used memory (GB)",
          "",
          "Message broker status",
          "  RabbitMQ reachability",
          "  Queue depth per worker",
          "",
          "Celery worker fleet",
          "  Active workers listed",
          "  Tasks per queue type",
          "",
          "Instant alert if any",
          "component goes offline"]),

        (BLUE,   FILL_BLUE,   "ING", "Ingestion Success Rate",
         ["Real-time job funnel:",
          "  Waiting  (pending metadata)",
          "  Processing (active tasks)",
          "  Ready     (successful)",
          "  Error     (failed)",
          "",
          "Success rate =",
          "  ready / (ready + error)",
          "",
          "Avg processing time",
          "  seconds per document",
          "",
          "Top 5 error messages",
          "  for rapid debugging"]),

        (GREEN,  FILL_GREEN,  "RAG", "RAG Quality Score",
         ["Per-response quality labels",
          "scored by eval pipeline:",
          "  HIGH   — full context hit",
          "  MEDIUM — partial context",
          "  LOW    — poor retrieval",
          "",
          "Semantic cache hit rate",
          "  (Redis answer caching)",
          "",
          "Avg response latency (ms)",
          "",
          "Worst 5 documents by",
          "low-quality rate flagged",
          "for re-ingestion review"]),

        (PURPLE, FILL_PURPLE, "FB",  "User Feedback Loop",
         ["Users rate every response:",
          "  Correct / Helpful",
          "  Incorrect / Hallucination",
          "  Missing context",
          "",
          "Live feedback counters:",
          "  Total feedback events",
          "  Positive vs Negative",
          "  Average feedback score",
          "",
          "Label distribution chart",
          "  Top 8 labels ranked",
          "",
          "Feeds model retraining",
          "and prompt improvement"]),
    ]

    for i, (ac, fi, ico, ttl, buls) in enumerate(cards):
        card(slide, 0.20 + i * (cw + gap), cy, cw, ch, ac, fi, ttl, buls, ico)


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    global _LOGO

    print(f"Opening: {PPT_PATH}")
    prs = Presentation(str(PPT_PATH))
    n   = len(prs.slides)
    print(f"  Current slide count: {n}")

    # ── Step 1: Remove previously added slides (31-41) if present ──────────────
    if n > 30:
        print(f"  Removing {n - 30} previously added slides...")
        remove_slides_from(prs, 30)

    print(f"  Clean slide count: {len(prs.slides)}")

    # ── Step 2: Extract KAVIN logo from slide 1 ────────────────────────────────
    try:
        for shape in list(prs.slides)[0].shapes:
            if shape.shape_type == 13:   # PICTURE
                _LOGO = shape.image.blob
        print(f"  KAVIN logo extracted: {len(_LOGO):,} bytes")
    except Exception as e:
        print(f"  Logo extract failed: {e}")

    # ── Step 3: Add 2 new slides at the END first (then move into position) ────
    INSERT_AFTER = 9   # after slide 9 (Agentic Review Pipeline), 0-indexed = 8

    print(f"\n  Building 2 slides and inserting after slide {INSERT_AFTER}...")

    blank = prs.slide_layouts[6]   # Blank layout

    # Build Slide A: Admin Dashboard
    slA = prs.slides.add_slide(blank)
    build_admin_dashboard(slA)
    # Move to position INSERT_AFTER (0-based)
    move_slide_to(prs, len(prs.slides) - 1, INSERT_AFTER)
    print(f"  Inserted 'Developer Admin Dashboard' at position {INSERT_AFTER + 1}")

    # Build Slide B: Analytics
    slB = prs.slides.add_slide(blank)
    build_analytics(slB)
    # Move to position INSERT_AFTER + 1
    move_slide_to(prs, len(prs.slides) - 1, INSERT_AFTER + 1)
    print(f"  Inserted 'System Quality & Analytics' at position {INSERT_AFTER + 2}")

    print(f"\n  Total slides: {len(prs.slides)}")

    # Verify ordering
    print("\n  Slide order verify (slides 8-13):")
    for i, slide in enumerate(list(prs.slides)[7:13], 8):
        t = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type != 13:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    break
        print(f"    Slide {i:2d}: {t}")

    # ── Step 4: Save ───────────────────────────────────────────────────────────
    try:
        prs.save(str(PPT_PATH))
        print(f"\n  Saved -> {PPT_PATH}")
    except PermissionError:
        alt = PPT_PATH.parent / "KavinBase_Final_v2.pptx"
        prs.save(str(alt))
        print(f"\n  File locked — saved -> {alt}")


if __name__ == "__main__":
    main()
