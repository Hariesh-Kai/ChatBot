"""
add_devtools_slides_v2.py
Appends Developer-Mode slides to KavinBase_Presentation_Updated.pptx
using the exact same dark-navy / blue-header / card design of the existing slides.
"""

import io, base64, zlib, requests
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR   = Path(__file__).resolve().parent
INPUT_PPT  = BASE_DIR / "KavinBase_Presentation_Updated.pptx"
OUTPUT_PPT = BASE_DIR / "KavinBase_Presentation_DevTools_v2.pptx"

# ── Colour Palette (matches existing slides) ─────────────────────────────────
BG        = RGBColor(0x05, 0x0A, 0x18)   # near-black navy
HEADER_BG = RGBColor(0x1E, 0x4A, 0xC4)   # royal blue header
CARD_BG   = RGBColor(0x0D, 0x1B, 0x38)   # card fill
LOGO_BG   = RGBColor(0x0A, 0x14, 0x2E)   # logo box fill
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUB_COL   = RGBColor(0x93, 0xC5, 0xFD)   # light-blue subtitle
BODY_COL  = RGBColor(0xD8, 0xE8, 0xF8)   # body text
DIM_COL   = RGBColor(0x6A, 0x8A, 0xB0)   # dim / arrow text

BLUE   = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
RED    = RGBColor(0xEF, 0x44, 0x44)
CYAN   = RGBColor(0x06, 0xB6, 0xD4)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
YELLOW = RGBColor(0xFB, 0xBC, 0x04)

FONT = "Calibri"

# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _run(para, text, size=13, color=None, bold=False, italic=False):
    """Append a run to a paragraph."""
    from pptx.oxml.ns import qn
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color:
        run.font.color.rgb = color
    return run


def set_bg(slide):
    """Fill slide background with BG colour."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def add_rect(slide, x, y, w, h, fill=CARD_BG, border=BLUE, border_pt=1.5, no_border=False):
    """Add a filled (and optionally bordered) rectangle."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if no_border:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width    = Pt(border_pt)
    shape.shadow.inherit = False
    return shape


def add_tb(slide, x, y, w, h):
    """Add a transparent text-box."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def set_para(tf, idx, text, size=13, color=None, bold=False, italic=False,
             align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    """Set / add a paragraph on a text frame."""
    if idx == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.name   = FONT
    if color:
        r.font.color.rgb = color
    return p, r


# ── Chrome: header bar + KAVIN logo ──────────────────────────────────────────

def add_chrome(slide, title, subtitle, SW):
    """
    Add the blue header bar, subtitle, and KAVIN logo that appear on every slide.
    SW = slide width in inches.
    """
    logo_w   = max(1.6, SW * 0.135)
    logo_x   = SW - logo_w - 0.16
    hdr_w    = logo_x - 0.20
    hdr_x    = 0.18
    hdr_y    = 0.56
    hdr_h    = 0.78

    # Blue header bar
    hdr = add_rect(slide, hdr_x, hdr_y, hdr_w, hdr_h, fill=HEADER_BG, no_border=True)
    tf  = hdr.text_frame
    tf.margin_left   = Pt(14)
    tf.margin_top    = Pt(10)
    tf.margin_bottom = Pt(4)
    tf.word_wrap     = True
    p, r = set_para(tf, 0, title, size=22, color=WHITE, bold=True)

    # KAVIN logo box
    logo = add_rect(slide, logo_x, hdr_y, logo_w, hdr_h, fill=LOGO_BG, border=BLUE, border_pt=1.5)
    ltf  = logo.text_frame
    ltf.margin_top  = Pt(8)
    ltf.margin_left = Pt(6)
    lp, lr = set_para(ltf, 0, "  KAVIN", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    lp2, lr2 = set_para(ltf, 1, "Developer Tools", size=9, color=SUB_COL, align=PP_ALIGN.CENTER)

    # Subtitle
    sub_tb = add_tb(slide, hdr_x, hdr_y + hdr_h + 0.06, SW - 0.36, 0.32)
    sub_tf = sub_tb.text_frame
    sub_tf.margin_left = Pt(4)
    sub_tf.margin_top  = Pt(0)
    set_para(sub_tf, 0, subtitle, size=11, color=SUB_COL, italic=True)

    return hdr_y + hdr_h + 0.44   # returns Y where content should start


# ── Card helpers ──────────────────────────────────────────────────────────────

def fill_card(slide, x, y, w, h, accent, icon_label, title, bullets):
    """
    Draw a dark card with a coloured border, icon+title header row, and bullet list.
    """
    # Outer card rectangle
    card = add_rect(slide, x, y, w, h, fill=CARD_BG, border=accent, border_pt=1.8)

    tf = card.text_frame
    tf.word_wrap     = True
    tf.margin_left   = Pt(10)
    tf.margin_right  = Pt(6)
    tf.margin_top    = Pt(10)
    tf.margin_bottom = Pt(6)

    # Icon + title (first paragraph)
    p0, r0 = set_para(tf, 0, f"{icon_label}  {title}", size=14,
                      color=accent, bold=True, space_after=4)

    # Thin separator line (use a horizontal rectangle)
    sep = add_rect(slide, x + 0.08, y + 0.44, w - 0.16, 0.014,
                   fill=accent, no_border=True)

    # Bullet points (each is a separate paragraph)
    tb_x = x + 0.08
    tb_y = y + 0.48
    tb_w = w - 0.16
    tb_h = h - 0.56

    btb = add_tb(slide, tb_x, tb_y, tb_w, tb_h)
    btf = btb.text_frame
    btf.word_wrap     = True
    btf.margin_left   = Pt(4)
    btf.margin_top    = Pt(4)
    btf.margin_right  = Pt(4)
    btf.margin_bottom = Pt(4)

    for i, line in enumerate(bullets):
        sz = 12 if not line.startswith("  ") else 11
        col = BODY_COL if not line.startswith("  ") else DIM_COL
        set_para(btf, i, line, size=sz, color=col, space_before=0, space_after=1)

    return card


def fill_card_no_tf(slide, x, y, w, h, accent, icon_label, title, bullets):
    """Alias kept for compatibility — delegates to fill_card."""
    return fill_card(slide, x, y, w, h, accent, icon_label, title, bullets)


def add_process_steps(slide, sy, sh, SW, steps, colors):
    """
    Draw a horizontal sequence of labelled boxes with arrows between them.
    steps = [{"title": str, "body": str}, ...]
    """
    n       = len(steps)
    margin  = 0.20
    arr_w   = 0.28
    total_w = SW - margin * 2
    box_w   = (total_w - arr_w * (n - 1)) / n
    x       = margin

    for i, step in enumerate(steps):
        col = colors[i % len(colors)]
        box = add_rect(slide, x, sy, box_w, sh, fill=CARD_BG, border=col, border_pt=1.5)

        tf = box.text_frame
        tf.word_wrap     = True
        tf.margin_left   = Pt(6)
        tf.margin_right  = Pt(4)
        tf.margin_top    = Pt(8)
        tf.margin_bottom = Pt(4)

        set_para(tf, 0, step["title"], size=12, color=col,
                 bold=True, align=PP_ALIGN.CENTER, space_after=3)

        if step.get("body"):
            for j, line in enumerate(step["body"]):
                set_para(tf, j + 1, line, size=10, color=BODY_COL,
                         align=PP_ALIGN.CENTER)

        x += box_w

        if i < n - 1:
            # Arrow
            atb = add_tb(slide, x + 0.03, sy + sh / 2 - 0.18, arr_w - 0.06, 0.36)
            atf = atb.text_frame
            atf.margin_top = Pt(0)
            p, r = set_para(atf, 0, "\u2192", size=18, color=BLUE,
                            align=PP_ALIGN.CENTER)
            x += arr_w


def add_metric_box(slide, x, y, w, h, accent, label, value, desc):
    """Small metric / KPI box."""
    box = add_rect(slide, x, y, w, h, fill=CARD_BG, border=accent, border_pt=1.8)
    tb  = add_tb(slide, x + 0.08, y + 0.08, w - 0.16, h - 0.16)
    tf  = tb.text_frame
    tf.word_wrap = True
    tf.margin_top  = Pt(0)
    tf.margin_left = Pt(0)
    set_para(tf, 0, label, size=10, color=accent,  bold=True, space_after=2)
    set_para(tf, 1, value, size=20, color=WHITE,    bold=True, space_after=2)
    set_para(tf, 2, desc,  size=10, color=BODY_COL, italic=True)


def get_mermaid(text):
    try:
        compressed = zlib.compress(text.encode("utf-8"), 9)
        b64 = base64.urlsafe_b64encode(compressed).decode("ascii")
        url = f"https://kroki.io/mermaid/png/{b64}"
        r = requests.get(url, timeout=18)
        return io.BytesIO(r.content) if r.status_code == 200 else None
    except Exception:
        return None


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # Blank layout


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_devtools_overview(prs, SW, SH):
    """SLIDE A — Developer Mode Overview: 2×2 card grid."""
    print("  Building: Developer Mode Overview")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "Developer Mode  —  Overview",
                    "Admin-gated REST panel for deep system introspection  |  router prefix: /devtools  |  All endpoints require admin JWT", SW)

    gx = 0.20; gy = cy; gw = (SW - 0.60) / 2; gh = (SH - gy - 0.30) / 2; gap = 0.20

    cards = [
        (BLUE,   "AI",  "Model Management",
         ["Install, download & register LLMs at runtime",
          "Three modes: base (HF) / lite (GGUF) / net (Cloud)",
          "Test model output via POST /devtools/models/test",
          "Delete & clean up stale model files from disk",
          "Patch live model registry without restart"]),
        (PURPLE, "DEBUG", "RAG Debug Tools",
         ["POST /devtools/intent   — classify query intent",
          "POST /devtools/rewrite  — inspect query rewriting",
          "POST /devtools/keywords — keyword extraction view",
          "POST /devtools/retrieve — full retrieval dry-run",
          "Session state inspector per chat session"]),
        (GREEN,  "CHART", "Analytics Dashboard",
         ["GET /devtools/mvp/overview — command centre",
          "Ingestion funnel: waiting / processing / ready / error",
          "RAG quality: high / medium / low + cache hit %",
          "User feedback score & label distribution",
          "Runtime: GPU VRAM, RabbitMQ, worker queues"]),
        (RED,    "CTRL",  "System Controls",
         ["Full user lifecycle: create / disable / reset PW",
          "Per-user & per-session RAG override switches",
          "Live DB browser: pgvector, chat DB, Redis, MinIO",
          "Selective reset: RAG / Chat / Redis / MinIO / ALL",
          "Safety gate: env flag + confirmation phrase"]),
    ]
    for i, (col, icon, ttl, buls) in enumerate(cards):
        col_i = i % 2; row_i = i // 2
        cx = gx + col_i * (gw + gap)
        cy2 = gy + row_i * (gh + gap * 0.6)
        fill_card(sl, cx, cy2, gw, gh, col, icon, ttl, buls)


def slide_model_management(prs, SW, SH):
    """SLIDE B — LLM Model Management: 3-column mode cards + optional flowchart."""
    print("  Building: LLM Model Management")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "LLM Model Management  —  Three Runtime Modes",
                    "Download, register, switch and test language models at runtime without server restart", SW)

    cw = (SW - 0.60) / 3; ch = SH - cy - 0.30; gap = 0.20

    cards = [
        (BLUE,   "BASE", "base  Mode  —  HuggingFace",
         ["Full-precision transformer models",
          "GPU (CUDA) preferred, CPU fallback",
          "Loaded into memory via HF pipeline",
          "",
          "Known models:",
          "  Qwen/Qwen2.5-3B-Instruct",
          "  Qwen/Qwen2.5-7B-Instruct",
          "",
          "Endpoint: POST /devtools/models/hf/install",
          "Caches to: models/hf/ (HF snapshot dir)",
          "",
          "Best for: high-accuracy answers on GPU servers"]),
        (PURPLE, "LITE", "lite  Mode  —  GGUF / llama.cpp",
         ["4-bit quantised GGUF models",
          "Runs entirely on CPU (RAM, no GPU needed)",
          "Loaded via llama-cpp-python bindings",
          "",
          "Known models:",
          "  Qwen2.5-3B-Instruct-GGUF  (Q4)",
          "  Qwen2.5-1.5B-Instruct-GGUF (Q4)",
          "  Qwen2.5-7B-Instruct-GGUF  (Q4)",
          "  Meta-Llama-3.1-8B-GGUF",
          "",
          "Endpoint: POST /devtools/models/download",
          "Smart-detects GGUF repo, stores in models/gguf/"]),
        (GREEN,  "NET",  "net  Mode  —  Cloud API",
         ["Routes queries to external LLM providers",
          "No local model download required",
          "API key stored in backend secrets store",
          "",
          "Supported providers:",
          "  OpenAI  (GPT-4o, GPT-4-turbo)",
          "  Google  (Gemini Pro / Flash)",
          "  Anthropic, Cohere (extensible)",
          "",
          "Endpoint: PATCH /devtools/models/registry",
          "Key check: GET /devtools/models/active",
          "",
          "Best for: cloud-first deployments"]),
    ]
    for i, (col, icon, ttl, buls) in enumerate(cards):
        fill_card(sl, 0.20 + i * (cw + gap), cy, cw, ch, col, icon, ttl, buls)


def slide_model_download(prs, SW, SH):
    """SLIDE C — Model Download & Registry API: 2-column card."""
    print("  Building: Model Download & Registry")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "Model Download & Registry API",
                    "REST endpoints to install, register, test and remove models from the live backend", SW)

    cw = (SW - 0.60) / 2; ch = SH - cy - 0.30; gap = 0.20

    fill_card(sl, 0.20, cy, cw, ch, BLUE, "DL", "Download Endpoints",
        ["POST /devtools/models/download",
         "  Auto-detects GGUF vs HF snapshot repo",
         "  Handles stale / duplicate registrations",
         "",
         "POST /devtools/models/hf/install",
         "  snapshot_download to HF cache dir",
         "  Registers in model config automatically",
         "",
         "POST /devtools/models/gguf/download",
         "  Direct URL stream -> models/gguf/",
         "  25 GB safety cap, .part temp file",
         "",
         "POST /devtools/models/gguf/register",
         "  Register a pre-existing local .gguf path",
         "",
         "POST /devtools/models/test",
         "  Generate a sample to confirm model works",
         "",
         "DELETE /devtools/models/{model_id}",
         "  Unregisters and deletes files from disk"])

    fill_card(sl, 0.20 + cw + gap, cy, cw, ch, PURPLE, "AUX", "Auxiliary & Embedding Models",
        ["EMBEDDING",
         "  BAAI/bge-m3  ->  384-dim dense vectors",
         "  Used by PGVector for semantic search",
         "",
         "INTENT CLASSIFICATION",
         "  facebook/bart-large-mnli",
         "  Zero-shot NLI: fact_lookup / comparison",
         "  /comparison / procedural ...",
         "",
         "LAYOUT / OCR",
         "  microsoft/table-transformer",
         "  unstructuredio/yolo_x_layout",
         "  Used during hi_res PDF preprocessing",
         "",
         "REGISTRY PATCH",
         "  PATCH /devtools/models/registry",
         "  JSON body: {lite:{default:..}, base:{..}}",
         "  Hot-swaps active model without restart",
         "",
         "GET  /devtools/models         (inventory)",
         "GET  /devtools/models/active  (readiness)"])


def slide_rag_debug(prs, SW, SH):
    """SLIDE D — RAG Debug & Retrieval Testing: process flow + cards."""
    print("  Building: RAG Debug & Retrieval Testing")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "RAG Debug & Retrieval Testing",
                    "Step-by-step pipeline testing tools  |  Each stage can be called independently via a dedicated devtools endpoint", SW)

    # Process flow row
    steps = [
        {"title": "1. Input",           "body": ["Raw developer", "question string"]},
        {"title": "2. normalize_text",  "body": ["Lowercase, strip", "whitespace, unicode"]},
        {"title": "3. Intent Classify", "body": ["bart-large-mnli", "fact_lookup / compare"]},
        {"title": "4. Query Rewrite",   "body": ["Optional: expand", "with chat history"]},
        {"title": "5. Keyword Extract", "body": ["SQL BM25 terms", "for hybrid search"]},
        {"title": "6. PGVector Fetch",  "body": ["Dense semantic", "retrieval from DB"]},
        {"title": "7. Top N Chunks",    "body": ["Ranked results", "returned to caller"]},
    ]
    flow_h = 1.65
    add_process_steps(sl, cy, flow_h, SW, steps,
                      [BLUE, CYAN, PURPLE, GREEN, ORANGE, RED, BLUE])

    # Card row below
    card_y = cy + flow_h + 0.28
    card_h = SH - card_y - 0.28
    cw = (SW - 0.60) / 4; gap = 0.20

    fill_card(sl, 0.20,            card_y, cw, card_h, BLUE,   "INT",  "/devtools/intent",
        ["Classify any text string",
         "Returns normalized form",
         "and detected intent label",
         "Intents: fact_lookup",
         "  comparison, procedural",
         "  definition, multi_doc"])
    fill_card(sl, 0.20+cw+gap,     card_y, cw, card_h, PURPLE, "RW",   "/devtools/rewrite",
        ["Show query rewriting",
         "Pass mock chat history",
         "Returns original +",
         "rewritten question",
         "Uses LLM-based",
         "context expansion"])
    fill_card(sl, 0.20+2*(cw+gap), card_y, cw, card_h, GREEN,  "KW",   "/devtools/keywords",
        ["Inspect keyword tokens",
         "extracted for SQL search",
         "Used in hybrid BM25+",
         "semantic mode",
         "Helps debug why a doc",
         "is / isn't retrieved"])
    fill_card(sl, 0.20+3*(cw+gap), card_y, cw, card_h, CYAN,   "RTR",  "/devtools/retrieve",
        ["Full pipeline dry-run",
         "Pass: question, doc ID,",
         "revision, collection",
         "Returns top-N chunks,",
         "chunk_ids, rag_mode,",
         "intent, and preview[0:3]"])


def slide_mvp_overview(prs, SW, SH):
    """SLIDE E — MVP Analytics Command Center: 4 metric cards + info."""
    print("  Building: MVP Analytics Command Center")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "MVP Analytics Command Center",
                    "GET /devtools/mvp/overview  |  window_hours: 1-168  |  Admin-gated real-time operational snapshot", SW)

    cw = (SW - 0.60) / 4; ch = SH - cy - 0.30; gap = 0.20
    cards = [
        (RED,    "RT",   "Runtime Status",
         ["GPU: device, VRAM total/free",
          "CUDA version + torch build",
          "",
          "RabbitMQ broker:",
          "  reachability ping",
          "  queue depth per worker",
          "",
          "Celery workers:",
          "  active worker list",
          "  task count per queue",
          "",
          "Software versions:",
          "  Python, torch, llama-cpp"]),
        (BLUE,   "INGEST","Ingestion Funnel",
         ["Per-window job counts:",
          "  WAIT_FOR_METADATA",
          "  PROCESSING (active)",
          "  READY  (success)",
          "  ERROR  (failed)",
          "  TOTAL  (all statuses)",
          "",
          "Success rate",
          "  = ready / (ready+error)",
          "",
          "Avg completion seconds",
          "  (READY + ERROR jobs)",
          "",
          "Top 5 error messages"]),
        (GREEN,  "RAG",  "RAG Quality Metrics",
         ["Total Q&A turns in window",
          "",
          "Quality distribution:",
          "  HIGH  quality responses",
          "  MEDIUM quality",
          "  LOW   quality",
          "  UNKNOWN (unscored)",
          "",
          "Cache hit rate",
          "  (Redis semantic cache)",
          "",
          "Avg latency  (ms)",
          "",
          "Worst 5 docs by",
          "  low-quality rate"]),
        (PURPLE, "FB",   "User Feedback",
         ["Total feedback events",
          "",
          "POSITIVE labels:",
          "  correct / helpful",
          "  thumbs_up",
          "",
          "NEGATIVE labels:",
          "  incorrect",
          "  hallucination",
          "  missing_context",
          "",
          "Average feedback score",
          "",
          "Label distribution",
          "  (top 8 labels ranked)"]),
    ]
    for i, (col, icon, ttl, buls) in enumerate(cards):
        fill_card(sl, 0.20 + i * (cw + gap), cy, cw, ch, col, icon, ttl, buls)


def slide_upload_inspector(prs, SW, SH):
    """SLIDE F — Upload Job Inspector: 2-column card."""
    print("  Building: Upload Job Inspector")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "Upload Job Inspector & Preprocessing Preview",
                    "GET /devtools/uploads/recent  |  Last N ingestion jobs with full metadata, artifact links and preview flags", SW)

    cw = (SW - 0.60) / 2; ch = SH - cy - 0.30; gap = 0.20

    fill_card(sl, 0.20, cy, cw, ch, BLUE, "JOB", "Job Record Fields",
        ["job_id          unique Celery task UUID",
         "session_id      originating chat session",
         "status          PROCESSING / READY / ERROR",
         "progress        0-100 integer percentage",
         "progress_label  human-readable stage name",
         "error           last error message if any",
         "",
         "company_document_id  doc reference code",
         "revision_number       doc version tag",
         "source_file           original filename",
         "",
         "rag_preprocessor   pypdf / pymupdf / unstructured",
         "rag_ingest_mode   fast / balanced / hi_fi",
         "rag_collection_name  PGVector collection used",
         "",
         "created_at / updated_at  timestamps"])

    fill_card(sl, 0.20 + cw + gap, cy, cw, ch, GREEN, "PRV", "Preview & Artifact Flags",
        ["preview_available",
         "  True if PDF or cached artifacts exist",
         "",
         "artifact_only_preview",
         "  PDF gone but cached JSON artifacts remain",
         "  Artifacts: raw_elements.json",
         "             filtered_elements.json",
         "             removed_elements.json",
         "             filter_report.json",
         "             chunks.json",
         "             enriched_chunks.json",
         "",
         "segregation_summary",
         "  element_groups from filter_report.json",
         "  Breakdown: text / table / image blocks",
         "",
         "missing_fields",
         "  List of required metadata not yet provided",
         "  (WAIT_FOR_METADATA status trigger)"])


def slide_user_management(prs, SW, SH):
    """SLIDE G — User Management & Provisioning: 3-column card."""
    print("  Building: User Management & Provisioning")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "User Management & Resource Provisioning",
                    "Admin-only lifecycle management  |  Every user gets isolated PostgreSQL DB + MinIO bucket + Redis namespace", SW)

    cw = (SW - 0.60) / 3; ch = SH - cy - 0.30; gap = 0.20

    fill_card(sl, 0.20, cy, cw, ch, BLUE, "NEW", "Create User",
        ["POST /devtools/users",
         "",
         "Payload fields:",
         "  email        unique email address",
         "  username     login handle",
         "  password     plain text (hashed)",
         "  role         user | admin",
         "  pg_database  optional, auto-derived",
         "  minio_bucket optional, auto-derived",
         "",
         "Auto-provisions:",
         "  PostgreSQL DB  chat_ui_{username}",
         "  MinIO bucket   chat-ui-{username}",
         "  Redis namespace user:{username}:",
         "",
         "Rollback: deletes user if provisioning fails"])

    fill_card(sl, 0.20 + cw + gap, cy, cw, ch, PURPLE, "MGR", "Manage Existing Users",
        ["GET /devtools/users",
         "  List all registered users",
         "",
         "PATCH /devtools/users/disable",
         "  Enable or disable login",
         "  Body: {identifier, disabled:bool}",
         "",
         "PATCH /devtools/users/password",
         "  Admin reset (no old pwd needed)",
         "  Body: {identifier, new_password}",
         "",
         "PATCH /devtools/users/role",
         "  Promote / demote to admin",
         "  Body: {identifier, role}",
         "",
         "DELETE /devtools/users",
         "POST   /devtools/users/delete",
         "  Remove user record from auth store"])

    fill_card(sl, 0.20 + 2*(cw+gap), cy, cw, ch, GREEN, "RES", "Resource Provisioning Detail",
        ["PostgreSQL Database",
         "  _ensure_postgres_database()",
         "  Creates DB if not exists",
         "  Runs _init_db() for chat schema",
         "  Name: 3-63 chars, auto-sanitised",
         "",
         "MinIO Object Bucket",
         "  _ensure_minio_bucket()",
         "  Creates bucket if not present",
         "  Name: lowercase, 3-63 chars",
         "",
         "Redis Namespace",
         "  Pattern: user:{username}:",
         "  Namespaced key isolation",
         "",
         "set_user_resources()",
         "  Persists provisioning metadata",
         "  to auth store for the user record"])


def slide_rag_overrides(prs, SW, SH):
    """SLIDE H — RAG Overrides & Session State: 2-column."""
    print("  Building: RAG Overrides & Session State")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "RAG Overrides, Session State & Feature Flags",
                    "Per-session and per-user toggles for retrieval behaviour  |  Live feature flag tuning without restart", SW)

    cw120 = (SW - 0.60) * 0.52; cw2 = SW - 0.60 - cw120; ch = SH - cy - 0.30; gap = 0.20

    fill_card(sl, 0.20, cy, cw120, ch, BLUE, "OVR", "RAG Override Controls",
        ["GET /devtools/rag/overrides",
         "  List all active session/user overrides",
         "",
         "POST /devtools/rag/disable",
         "  Body: {session_id?, username?}",
         "  Forces direct LLM — no RAG retrieval",
         "",
         "POST /devtools/rag/enable",
         "  Re-enables RAG for that session/user",
         "",
         "Use cases:",
         "  Debug hallucination vs. retrieval gap",
         "  Force base model answer for comparison",
         "  Temporarily bypass RAG for a power user",
         "",
         "Session State Inspector",
         "GET /devtools/session-state/{session_id}",
         "  postgres_message_count",
         "  recent_user_messages  (last 3)",
         "  active_topic          (Redis value)",
         "  used_chunk_ids_count  (Redis set size)"])

    fill_card(sl, 0.20 + cw120 + gap, cy, cw2, ch, PURPLE, "FLAGS", "Feature Flags",
        ["GET  /devtools/settings",
         "PATCH /devtools/settings",
         "",
         "enable_query_rewrite",
         "  Rewrites question with history",
         "",
         "enable_hybrid_retrieval",
         "  Dense + BM25 keyword search",
         "",
         "force_detailed_retrieval",
         "  Ignores mode profile limits",
         "",
         "rag_retrieval_mode",
         "  auto / semantic / keyword",
         "  / hybrid / detailed",
         "",
         "rag_collection_name",
         "  Switch target PGVector",
         "  collection on-the-fly",
         "",
         "Job Inspector",
         "GET /devtools/jobs",
         "  Active in-memory job states"])


def slide_db_visibility(prs, SW, SH):
    """SLIDE I — Database Visibility Panel: 4 cards."""
    print("  Building: Database Visibility Panel")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "Database Visibility Panel",
                    "Admin-gated live inspection of all persistence layers  |  GET /devtools/dbs/{db_id}/records  with pagination", SW)

    cw = (SW - 0.60) / 4; ch = SH - cy - 0.30; gap = 0.20
    cards = [
        (BLUE,   "PGV",  "rag_db  (pgvector)",
         ["PostgreSQL + pgvector",
          "",
          "Tables:",
          "  langchain_pg_embedding",
          "  langchain_pg_collection",
          "",
          "Records browser:",
          "  Paginated (limit/offset)",
          "  embedding column auto-",
          "  excluded (too large)",
          "",
          "Used for:",
          "  Semantic search index",
          "  384-dim BGE-M3 vectors",
          "  Per-doc, per-revision"]),
        (PURPLE, "CHAT", "chat_db  (postgres)",
         ["PostgreSQL chat memory",
          "",
          "Tables:",
          "  chat_messages",
          "  chat_sessions",
          "  session_topic_hints",
          "  session_active_documents",
          "  rag_job_runs",
          "  rag_audit_log",
          "  retrieval_feedback",
          "",
          "Stores:",
          "  Full conversation history",
          "  RAG quality audit events",
          "  User feedback records"]),
        (GREEN,  "RDS",  "redis  (key-value)",
         ["Redis key browser",
          "",
          "Key types shown:",
          "  string  — raw value",
          "  list    — list[N] count",
          "  set     — set[N] count",
          "  zset    — zset[N] count",
          "  hash    — hash[N] count",
          "",
          "Namespaces:",
          "  rag:*     retrieval cache",
          "  abort:*   stream kill flags",
          "  user:*    session state",
          "",
          "Values truncated at 500 chars"]),
        (ORANGE, "S3",   "minio  (object store)",
         ["MinIO object browser",
          "",
          "Fields shown:",
          "  object      full path/key",
          "  size        bytes",
          "  last_modified  ISO datetime",
          "",
          "Default bucket:",
          "  chat-ui-documents  (env)",
          "",
          "Lists ALL objects",
          "recursively in bucket",
          "",
          "Paginated: limit / offset",
          "returns total object count"]),
    ]
    for i, (col, icon, ttl, buls) in enumerate(cards):
        fill_card(sl, 0.20 + i * (cw + gap), cy, cw, ch, col, icon, ttl, buls)


def slide_runtime_status(prs, SW, SH):
    """SLIDE J — Runtime Status: 2x2 card grid."""
    print("  Building: Runtime Status")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "Runtime Status  —  System Health Monitoring",
                    "GET /devtools/runtime  |  GET /devtools/models/active  |  Real-time hardware & service readiness", SW)

    gw = (SW - 0.60) / 2; gh = (SH - cy - 0.38) / 2; gap = 0.20

    data = [
        (BLUE,   "GPU",  "GPU & Compute",
         ["Device name (e.g. NVIDIA RTX 3090)",
          "VRAM total  /  free  /  used  (GB)",
          "CUDA version detected by torch",
          "torch version + build config",
          "",
          "Used to determine:",
          "  base vs lite mode auto-selection",
          "  Whether model fits in VRAM",
          "",
          "Fallback: CPU-only mode reported"]),
        (RED,    "MQ",   "RabbitMQ Broker",
         ["Broker reachability ping result",
          "If unreachable: error message shown",
          "",
          "Queue depths reported:",
          "  default        general tasks",
          "  rag_ingest     document pipeline",
          "  rag_high_mem   large PDF jobs",
          "",
          "Used by Celery for task routing",
          "High queue depth = ingestion backlog"]),
        (GREEN,  "WRK",  "Celery Workers",
         ["Active worker list with hostnames",
          "Task count per worker instance",
          "Queue assignments per worker",
          "",
          "Worker types:",
          "  default       — light async tasks",
          "  rag_ingest    — PDF chunking",
          "  rag_high_mem  — unstructured hi_res",
          "",
          "Zero workers = upload pipeline stalled"]),
        (PURPLE, "SW",   "Active Model Readiness",
         ["GET /devtools/models/active",
          "",
          "For each mode (base / lite / net):",
          "  configured_model_id",
          "  effective_model_id (resolved)",
          "  type: gguf | hf | net",
          "  path: local file or cache dir",
          "  ready: file/snapshot exists",
          "  loaded: currently in RAM",
          "  error: reason if not ready",
          "",
          "net mode also reports:",
          "  provider name + api key present"]),
    ]
    for i, (col, icon, ttl, buls) in enumerate(data):
        col_i = i % 2; row_i = i // 2
        cx = 0.20 + col_i * (gw + gap)
        cy2 = cy + row_i * (gh + gap * 0.55)
        fill_card(sl, cx, cy2, gw, gh, col, icon, ttl, buls)


def slide_system_reset(prs, SW, SH):
    """SLIDE K — System Reset Controls: warning card + 5 process steps."""
    print("  Building: System Reset Controls")
    sl = blank_slide(prs); set_bg(sl)
    cy = add_chrome(sl, "System Reset Controls  —  Admin Destructive Operations",
                    "Guarded by env flag + confirmation phrase  |  Each layer can be reset independently or all at once", SW)

    # Safety requirements card
    warn_h = 0.90
    warn = add_rect(sl, 0.20, cy, SW - 0.40, warn_h, fill=RGBColor(0x3A, 0x08, 0x08),
                    border=RED, border_pt=2.0)
    wtf = warn.text_frame
    wtf.margin_left = Pt(12); wtf.margin_top = Pt(8)
    wtf.word_wrap = True
    set_para(wtf, 0, "SAFETY REQUIREMENTS  —  Both conditions must be met before any reset fires:", size=12,
             color=RED, bold=True, space_after=3)
    set_para(wtf, 1,
             "  1.  Environment variable  CHAT_UI_ENABLE_DESTRUCTIVE_DEVTOOLS=1  must be set on the server.",
             size=11, color=BODY_COL)
    set_para(wtf, 2,
             "  2.  Request body must include:   confirm = 'DELETE_EVERYTHING'",
             size=11, color=WHITE, bold=True)

    # 5 reset steps as process boxes
    steps = [
        {"title": "POST /reset/rag",
         "body":  ["TRUNCATE", "langchain_pg_embedding", "langchain_pg_collection", "(clears all vectors)"]},
        {"title": "POST /reset/chat",
         "body":  ["TRUNCATE", "chat_messages", "chat_sessions", "session_topics + docs"]},
        {"title": "POST /reset/redis",
         "body":  ["Delete rag:*", "and abort:* keys", "or full flushdb", "if wipe_all=True"]},
        {"title": "POST /reset/minio",
         "body":  ["Remove ALL objects", "in target bucket", "minio_bucket field", "required in body"]},
        {"title": "POST /reset/all",
         "body":  ["Chains all 4", "resets in sequence", "Returns combined", "result per target"]},
    ]
    flow_y = cy + warn_h + 0.22
    flow_h = SH - flow_y - 0.28
    add_process_steps(sl, flow_y, flow_h, SW, steps,
                      [BLUE, PURPLE, GREEN, ORANGE, RED])


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    print(f"Opening: {INPUT_PPT}")
    prs = Presentation(str(INPUT_PPT))
    SW = prs.slide_width.inches
    SH = prs.slide_height.inches
    print(f"  Slide size: {SW:.2f}\" x {SH:.2f}\"")
    print(f"  Existing slides: {len(prs.slides)}")
    print()

    slide_devtools_overview(prs, SW, SH)
    slide_model_management(prs, SW, SH)
    slide_model_download(prs, SW, SH)
    slide_rag_debug(prs, SW, SH)
    slide_mvp_overview(prs, SW, SH)
    slide_upload_inspector(prs, SW, SH)
    slide_user_management(prs, SW, SH)
    slide_rag_overrides(prs, SW, SH)
    slide_db_visibility(prs, SW, SH)
    slide_runtime_status(prs, SW, SH)
    slide_system_reset(prs, SW, SH)

    total = len(prs.slides)
    print()
    try:
        prs.save(str(OUTPUT_PPT))
        print(f"Saved -> {OUTPUT_PPT}")
        print(f"Total slides: {total}")
    except Exception as e:
        print(f"Save failed: {e}")


if __name__ == "__main__":
    main()
