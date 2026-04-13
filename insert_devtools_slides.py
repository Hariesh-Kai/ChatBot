"""
insert_devtools_slides.py
Appends 11 Developer-Mode slides into KavinBase_Presentation_Updated.pptx
using the EXACT design of the existing slides (colors, fonts, header, logo).
"""

import io, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Force UTF-8 output so emojis in slide text don't crash
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

BASE_DIR   = Path(__file__).resolve().parent
PPT_PATH   = BASE_DIR / "KavinBase_Presentation_Updated.pptx"

# ── Exact palette extracted from existing slides ───────────────────────────────
BG      = RGBColor(0x12, 0x14, 0x1A)   # full-slide dark background
HEADER  = RGBColor(0x5B, 0x8D, 0xFF)   # blue header bar
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)   # title text
SUB     = RGBColor(0xDF, 0xE9, 0xFF)   # header subtitle text
BODY    = RGBColor(0xCC, 0xD6, 0xF1)   # general body text
DIM     = RGBColor(0x70, 0x80, 0xA0)   # dimmed/secondary text

BLUE    = RGBColor(0x5B, 0x8D, 0xFF)   # blue accent
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
GREEN   = RGBColor(0x4A, 0xDE, 0x80)   # green accent
RED     = RGBColor(0xFF, 0x45, 0x45)   # red accent
CYAN    = RGBColor(0x22, 0xD3, 0xEE)   # cyan accent
ORANGE  = RGBColor(0xFB, 0xBF, 0x24)   # orange accent

# Card fills (dark tinted per accent)
FILL_BLUE   = RGBColor(0x06, 0x12, 0x28)
FILL_PURPLE = RGBColor(0x10, 0x08, 0x28)
FILL_GREEN  = RGBColor(0x05, 0x18, 0x10)
FILL_RED    = RGBColor(0x1E, 0x06, 0x06)
FILL_CYAN   = RGBColor(0x04, 0x14, 0x20)
FILL_ORANGE = RGBColor(0x1E, 0x12, 0x04)
FILL_DARK   = RGBColor(0x0A, 0x10, 0x1E)

FONT    = "Calibri"
SW      = 13.33   # slide width inches
SH      = 7.50    # slide height inches
HDR_H   = 1.10    # header bar height

_KAVIN_LOGO_BYTES: bytes = b""   # populated at runtime


# ── Shape helpers ──────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=FILL_DARK, border=None, bpt=1.5):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def tb(slide, x, y, w, h):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.text_frame.word_wrap = True
    return t


def para(tf, text, size=13, color=BODY, bold=False, italic=False,
         align=PP_ALIGN.LEFT, idx=None, spc_before=0, spc_after=1):
    p = tf.paragraphs[0] if (idx == 0 or idx is None and len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(spc_before)
    p.space_after  = Pt(spc_after)
    r = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.name = FONT; r.font.color.rgb = color
    return p


# ── Slide chrome: background + header + logo ───────────────────────────────────

def chrome(slide, title, subtitle):
    """Add dark BG, blue header bar, title, subtitle, and KAVIN logo."""
    # 1. Dark background
    rect(slide, 0, 0, SW, SH, fill=BG)

    # 2. Blue header bar (full width, 1.1" tall)
    rect(slide, 0, 0, SW, HDR_H, fill=HEADER)

    # 3. Title in header
    t1 = tb(slide, 0.35, 0.10, SW - 2.5, 0.60)
    tf1 = t1.text_frame; tf1.margin_top = Pt(0); tf1.margin_left = Pt(4)
    para(tf1, title, size=26, color=WHITE, bold=True)

    # 4. Subtitle in header
    t2 = tb(slide, 0.35, 0.66, SW - 2.5, 0.36)
    tf2 = t2.text_frame; tf2.margin_top = Pt(0); tf2.margin_left = Pt(4)
    para(tf2, subtitle, size=13, color=SUB)

    # 5. KAVIN logo (picture if extracted, else text badge)
    if _KAVIN_LOGO_BYTES:
        slide.shapes.add_picture(
            io.BytesIO(_KAVIN_LOGO_BYTES),
            Inches(SW - 2.10), Inches(0.20),
            width=Inches(1.83), height=Inches(0.54)
        )
    else:
        logo = rect(slide, SW - 2.10, 0.16, 1.83, 0.60, fill=RGBColor(0x06, 0x12, 0x28), border=BLUE, bpt=1.2)
        ltf  = logo.text_frame; ltf.margin_top = Pt(6); ltf.margin_left = Pt(6)
        para(ltf, "KAVIN", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Returns Y where content should begin
    return HDR_H + 0.16


# ── Card helper ────────────────────────────────────────────────────────────────

def card(slide, x, y, w, h, accent, fill_color, title, lines, icon=""):
    # Outer border
    c = rect(slide, x, y, w, h, fill=fill_color, border=accent, bpt=1.8)

    # Title row
    title_h = 0.46
    title_tb = tb(slide, x + 0.10, y + 0.09, w - 0.20, title_h)
    ttf = title_tb.text_frame; ttf.margin_top = Pt(0); ttf.word_wrap = False
    lbl = f"{icon}  {title}" if icon else title
    para(ttf, lbl, size=14, color=accent, bold=True)

    # Thin separator line
    rect(slide, x + 0.08, y + title_h + 0.05, w - 0.16, 0.013, fill=accent)

    # Bullet body
    body_tb = tb(slide, x + 0.12, y + title_h + 0.10, w - 0.22, h - title_h - 0.15)
    btf = body_tb.text_frame; btf.word_wrap = True; btf.margin_top = Pt(2)
    for i, line in enumerate(lines):
        sz   = 11 if line.startswith("  ") else 12
        col  = DIM  if line.startswith("  ") else BODY
        para(btf, line, size=sz, color=col, spc_after=1)


# ── Process-step flow ─────────────────────────────────────────────────────────

def flow(slide, fy, fh, steps, colors, label_size=11):
    n = len(steps)
    arr_w = 0.26
    margin = 0.22
    total  = SW - margin * 2
    bw     = (total - arr_w * (n - 1)) / n
    x = margin

    for i, step in enumerate(steps):
        col   = colors[i % len(colors)]
        fills = [FILL_BLUE, FILL_PURPLE, FILL_GREEN, FILL_RED, FILL_CYAN, FILL_ORANGE, FILL_DARK]
        fi    = fills[i % len(fills)]
        b = rect(slide, x, fy, bw, fh, fill=fi, border=col, bpt=1.5)
        btf = b.text_frame; btf.word_wrap = True
        btf.margin_left = Pt(5); btf.margin_top = Pt(7); btf.margin_right = Pt(4)

        title_line = step.get("title", "")
        para(btf, title_line, size=label_size, color=col, bold=True,
             align=PP_ALIGN.CENTER, spc_after=3)

        for ln in step.get("body", []):
            para(btf, ln, size=label_size - 1, color=BODY, align=PP_ALIGN.CENTER, spc_after=1)

        x += bw
        if i < n - 1:
            at = tb(slide, x + 0.02, fy + fh / 2 - 0.18, arr_w - 0.04, 0.36)
            at.text_frame.margin_top = Pt(0)
            para(at.text_frame, "->", size=16, color=DIM, align=PP_ALIGN.CENTER)
            x += arr_w


# ── Individual slide builders ──────────────────────────────────────────────────

def s_overview(slide):
    cy = chrome(slide,
                "Developer Mode  -  Overview",
                "Admin-gated REST panel at /devtools  |  All endpoints require admin JWT  |  No restart needed for any change")
    gw = (SW - 0.58) / 2; gh = (SH - cy - 0.32) / 2; gap = 0.18

    cards_data = [
        (BLUE,   FILL_BLUE,   "AI",    "Model Management",
         ["Install, download & register LLMs at runtime",
          "3 modes: base (HF) / lite (GGUF) / net (Cloud)",
          "Test model output with a custom prompt",
          "Delete & clean up stale model files from disk",
          "Patch live model registry without restart"]),
        (PURPLE, FILL_PURPLE, "DBG",   "RAG Debug Tools",
         ["POST /devtools/intent    - classify query intent",
          "POST /devtools/rewrite   - inspect query rewriting",
          "POST /devtools/keywords  - keyword extraction view",
          "POST /devtools/retrieve  - full retrieval dry-run",
          "Session state inspector per chat session ID"]),
        (GREEN,  FILL_GREEN,  "CHART", "Analytics Dashboard",
         ["GET /devtools/mvp/overview  (1-168h window)",
          "Ingestion funnel: waiting / processing / ready / error",
          "RAG quality: high / medium / low + cache hit %",
          "User feedback: positive / negative / avg score",
          "Runtime: GPU VRAM, RabbitMQ, Celery workers"]),
        (RED,    FILL_RED,    "CTRL",  "System Controls",
         ["Full user lifecycle: create / disable / reset PW",
          "Per-user & per-session RAG override switches",
          "Live DB browser: pgvector, chat DB, Redis, MinIO",
          "Selective reset: RAG / Chat / Redis / MinIO / ALL",
          "Safety gate: env flag + confirmation phrase"]),
    ]
    for i, (ac, fi, ico, ttl, buls) in enumerate(cards_data):
        ci = i % 2; ri = i // 2
        card(slide, 0.20 + ci * (gw + gap), cy + ri * (gh + gap * 0.7),
             gw, gh, ac, fi, ttl, buls, ico)


def s_model_modes(slide):
    cy = chrome(slide,
                "LLM Model Management  -  Three Runtime Modes",
                "Switch between local HF transformers, CPU-friendly GGUF, and cloud API providers without a restart")
    cw = (SW - 0.58) / 3; ch = SH - cy - 0.28; gap = 0.18

    cards_data = [
        (BLUE, FILL_BLUE, "BASE", "base Mode  -  HuggingFace",
         ["Full-precision transformer models",
          "GPU (CUDA) preferred, CPU fallback",
          "Loaded via HF pipeline + tokenizer",
          "",
          "Known models:",
          "  Qwen/Qwen2.5-3B-Instruct",
          "  Qwen/Qwen2.5-7B-Instruct",
          "",
          "Install: POST /devtools/models/hf/install",
          "Cache:   models/hf/  (HF snapshot dir)",
          "",
          "Best for: high-accuracy on GPU servers"]),
        (PURPLE, FILL_PURPLE, "LITE", "lite Mode  -  GGUF / llama.cpp",
         ["4-bit quantised GGUF models",
          "Runs entirely on CPU (no GPU needed)",
          "Loaded via llama-cpp-python bindings",
          "",
          "Known models:",
          "  Qwen2.5-3B-Instruct-GGUF  (Q4)",
          "  Qwen2.5-1.5B-Instruct-GGUF (Q4)",
          "  Qwen2.5-7B-Instruct-GGUF  (Q4)",
          "  Meta-Llama-3.1-8B-GGUF",
          "",
          "Install: POST /devtools/models/download",
          "Store:   models/gguf/"]),
        (GREEN, FILL_GREEN, "NET", "net Mode  -  Cloud API",
         ["Routes queries to external providers",
          "No local model download required",
          "API key stored in backend secrets",
          "",
          "Supported providers:",
          "  OpenAI  (GPT-4o, GPT-4-turbo)",
          "  Google  (Gemini Pro / Flash)",
          "  Anthropic, Cohere (extensible)",
          "",
          "Key check: GET /devtools/models/active",
          "Patch:     PATCH /devtools/models/registry",
          "",
          "Best for: cloud-first deployments"]),
    ]
    for i, (ac, fi, ico, ttl, buls) in enumerate(cards_data):
        card(slide, 0.20 + i * (cw + gap), cy, cw, ch, ac, fi, ttl, buls, ico)


def s_model_api(slide):
    cy = chrome(slide,
                "Model Download & Registry API",
                "Install, register, test and remove models from the live backend via REST  |  Changes take effect immediately")
    cw = (SW - 0.58) / 2; ch = SH - cy - 0.28; gap = 0.18

    card(slide, 0.20, cy, cw, ch, BLUE, FILL_BLUE, "Download Endpoints",
         ["POST /devtools/models/download",
          "  Auto-detects GGUF vs HF snapshot",
          "  Handles stale / duplicate entries",
          "",
          "POST /devtools/models/hf/install",
          "  snapshot_download to HF cache dir",
          "",
          "POST /devtools/models/gguf/download",
          "  Direct URL stream to models/gguf/",
          "  25 GB safety cap, .part temp file",
          "",
          "POST /devtools/models/gguf/register",
          "  Register pre-existing local .gguf",
          "",
          "POST /devtools/models/test",
          "  Run a quick generation sample",
          "",
          "DELETE /devtools/models/{model_id}",
          "  Unregisters + deletes files from disk",
          "",
          "PATCH /devtools/models/registry",
          "  Patch lite/base/net model IDs live"], "DL")

    card(slide, 0.20 + cw + gap, cy, cw, ch, PURPLE, FILL_PURPLE, "Auxiliary & Embedding Models",
         ["EMBEDDING MODEL",
          "  BAAI/bge-m3  ->  384-dim dense vectors",
          "  Powers PGVector semantic search",
          "",
          "INTENT CLASSIFICATION",
          "  facebook/bart-large-mnli",
          "  Zero-shot NLI labeling:",
          "  fact_lookup / comparison / procedural",
          "",
          "LAYOUT / OCR (Unstructured)",
          "  microsoft/table-transformer",
          "  unstructuredio/yolo_x_layout",
          "  Used in hi_res PDF preprocessing",
          "",
          "MODEL INVENTORY",
          "  GET  /devtools/models",
          "         -> full inventory + registry",
          "  GET  /devtools/models/active",
          "         -> per-mode readiness check"], "AUX")


def s_rag_debug(slide):
    cy = chrome(slide,
                "RAG Debug & Retrieval Testing Pipeline",
                "Each stage can be tested independently  |  POST /devtools/intent  /rewrite  /keywords  /retrieve")

    # Process flow
    fh = 1.60
    steps = [
        {"title": "1. Input",          "body": ["raw question", "from developer"]},
        {"title": "2. normalize_text", "body": ["lowercase strip", "unicode clean"]},
        {"title": "3. Intent Classify","body": ["bart-large-mnli", "fact / compare"]},
        {"title": "4. Query Rewrite",  "body": ["optional LLM", "context expand"]},
        {"title": "5. Keyword Extract","body": ["BM25 tokens", "SQL hybrid"]},
        {"title": "6. PGVector Fetch", "body": ["dense semantic", "retrieval"]},
        {"title": "7. Top N Chunks",   "body": ["ranked results", "returned"]},
    ]
    flow(slide, cy, fh, steps,
         [BLUE, CYAN, PURPLE, GREEN, ORANGE, RED, BLUE])

    # 4 detail cards below
    card_y = cy + fh + 0.20
    card_h = SH - card_y - 0.28
    cw = (SW - 0.58) / 4; gap = 0.18

    card(slide, 0.20,             card_y, cw, card_h, BLUE,   FILL_BLUE,   "/devtools/intent",
         ["Classify any text string",
          "Returns normalized form",
          "and detected intent label",
          "Labels: fact_lookup",
          "  comparison, procedural",
          "  definition, multi_doc"], "INT")
    card(slide, 0.20+cw+gap,      card_y, cw, card_h, PURPLE, FILL_PURPLE, "/devtools/rewrite",
         ["Show query rewriting",
          "Pass mock chat history",
          "Returns original +",
          "rewritten question",
          "Uses LLM-based",
          "context expansion"], "RW")
    card(slide, 0.20+2*(cw+gap),  card_y, cw, card_h, GREEN,  FILL_GREEN,  "/devtools/keywords",
         ["Inspect keyword tokens",
          "extracted for SQL search",
          "Used in BM25 hybrid mode",
          "Helps debug why a doc",
          "is or is not retrieved",
          "from the vector store"], "KW")
    card(slide, 0.20+3*(cw+gap),  card_y, cw, card_h, CYAN,   FILL_CYAN,   "/devtools/retrieve",
         ["Full pipeline dry-run",
          "Body: question, doc ID,",
          "revision, collection",
          "Returns top-N chunks,",
          "chunk_ids, rag_mode,",
          "intent + preview[0:3]"], "RTR")


def s_mvp(slide):
    cy = chrome(slide,
                "MVP Analytics Command Center",
                "GET /devtools/mvp/overview  |  window_hours param: 1-168  |  Real-time operational snapshot for admins")
    cw = (SW - 0.58) / 4; ch = SH - cy - 0.28; gap = 0.18

    cards_data = [
        (RED,    FILL_RED,    "RT",    "Runtime Status",
         ["GPU device name",
          "VRAM total / free / used",
          "CUDA + torch versions",
          "",
          "RabbitMQ:",
          "  broker reachability",
          "  queue depth per worker",
          "",
          "Celery workers:",
          "  active worker list",
          "  task count per queue",
          "",
          "Python / llama-cpp build"]),
        (BLUE,   FILL_BLUE,   "INGEST","Ingestion Funnel",
         ["Per-window job counts:",
          "  WAIT_FOR_METADATA",
          "  PROCESSING (active)",
          "  READY       (success)",
          "  ERROR       (failed)",
          "  TOTAL       (all)",
          "",
          "Success rate:",
          "  ready / (ready+error)",
          "",
          "Avg completion seconds",
          "",
          "Top 5 error messages"]),
        (GREEN,  FILL_GREEN,  "RAG",   "RAG Quality",
         ["Total Q&A turns",
          "",
          "Quality distribution:",
          "  HIGH   responses",
          "  MEDIUM responses",
          "  LOW    responses",
          "  UNKNOWN (unscored)",
          "",
          "Cache hit rate (Redis)",
          "",
          "Avg latency (ms)",
          "",
          "Worst 5 docs by",
          "  low-quality rate"]),
        (PURPLE, FILL_PURPLE, "FB",    "User Feedback",
         ["Total feedback events",
          "",
          "POSITIVE labels:",
          "  correct / helpful",
          "  thumbs_up",
          "",
          "NEGATIVE labels:",
          "  incorrect",
          "  hallucination",
          "  missing_context",
          "",
          "Avg feedback score",
          "",
          "Top 8 label breakdown"]),
    ]
    for i, (ac, fi, ico, ttl, buls) in enumerate(cards_data):
        card(slide, 0.20 + i * (cw + gap), cy, cw, ch, ac, fi, ttl, buls, ico)


def s_upload(slide):
    cy = chrome(slide,
                "Upload Job Inspector & Preprocessing Preview",
                "GET /devtools/uploads/recent  |  Last N ingestion jobs with full metadata, artifacts and preview flags")
    cw = (SW - 0.58) / 2; ch = SH - cy - 0.28; gap = 0.18

    card(slide, 0.20, cy, cw, ch, BLUE, FILL_BLUE, "Job Record Fields",
         ["job_id          Celery task UUID",
          "session_id      originating session",
          "status          PROCESSING/READY/ERROR",
          "progress        0-100 percentage",
          "progress_label  current stage name",
          "error           last error if any",
          "",
          "company_document_id",
          "revision_number",
          "source_file  (original filename)",
          "",
          "rag_preprocessor",
          "  pypdf / pymupdf / unstructured",
          "rag_ingest_mode",
          "  fast / balanced / hi_fi",
          "rag_collection_name",
          "",
          "created_at  /  updated_at"], "JOB")

    card(slide, 0.20 + cw + gap, cy, cw, ch, GREEN, FILL_GREEN, "Preview & Artifact Flags",
         ["preview_available",
          "  True if PDF or artifacts exist",
          "",
          "artifact_only_preview",
          "  PDF gone, cached JSONs remain:",
          "  raw_elements.json",
          "  filtered_elements.json",
          "  removed_elements.json",
          "  filter_report.json",
          "  chunks.json",
          "  enriched_chunks.json",
          "",
          "segregation_summary",
          "  element_groups from filter report",
          "  text / table / image breakdown",
          "",
          "missing_fields",
          "  Required metadata not yet provided",
          "  (WAIT_FOR_METADATA trigger)"], "PRV")


def s_users(slide):
    cy = chrome(slide,
                "User Management & Resource Provisioning",
                "Admin-only lifecycle  |  Every user gets isolated PostgreSQL DB + MinIO bucket + Redis namespace")
    cw = (SW - 0.58) / 3; ch = SH - cy - 0.28; gap = 0.18

    card(slide, 0.20, cy, cw, ch, BLUE, FILL_BLUE, "Create User",
         ["POST /devtools/users",
          "",
          "Payload:",
          "  email       unique address",
          "  username    login handle",
          "  password    plain (hashed)",
          "  role        user | admin",
          "  pg_database  auto-derived",
          "  minio_bucket auto-derived",
          "",
          "Auto-provisions:",
          "  PostgreSQL:  chat_ui_{user}",
          "  MinIO:       chat-ui-{user}",
          "  Redis:       user:{user}:",
          "",
          "Rollback if provisioning fails"], "NEW")

    card(slide, 0.20 + cw + gap, cy, cw, ch, PURPLE, FILL_PURPLE, "Manage Users",
         ["GET /devtools/users",
          "  List all registered users",
          "",
          "PATCH /devtools/users/disable",
          "  Enable or disable login",
          "",
          "PATCH /devtools/users/password",
          "  Admin reset (no old pwd)",
          "",
          "PATCH /devtools/users/role",
          "  Promote / demote to admin",
          "",
          "DELETE /devtools/users",
          "POST   /devtools/users/delete",
          "  Both remove user record"], "MGR")

    card(slide, 0.20 + 2*(cw+gap), cy, cw, ch, GREEN, FILL_GREEN, "Resources Provisioned",
         ["PostgreSQL Database",
          "  _ensure_postgres_database()",
          "  Creates DB if not exists",
          "  Runs chat schema init",
          "  Name: 3-63 chars, sanitised",
          "",
          "MinIO Object Bucket",
          "  _ensure_minio_bucket()",
          "  Creates bucket if needed",
          "  Name: lowercase, 3-63 chars",
          "",
          "Redis Namespace",
          "  Pattern: user:{username}:",
          "  Key isolation per user",
          "",
          "set_user_resources()",
          "  Persists metadata to auth store"], "RES")


def s_overrides(slide):
    cy = chrome(slide,
                "RAG Overrides, Session State & Feature Flags",
                "Per-session and per-user retrieval toggles  |  GET/PATCH /devtools/settings  |  No restart required")
    cw1 = (SW - 0.58) * 0.52; cw2 = (SW - 0.58) - cw1; ch = SH - cy - 0.28; gap = 0.18

    card(slide, 0.20, cy, cw1, ch, BLUE, FILL_BLUE, "RAG Override Controls",
         ["GET /devtools/rag/overrides",
          "  List all active overrides",
          "",
          "POST /devtools/rag/disable",
          "  Body: {session_id?, username?}",
          "  Forces direct LLM - no RAG",
          "",
          "POST /devtools/rag/enable",
          "  Re-enables RAG for target",
          "",
          "Session State Inspector",
          "GET /devtools/session-state/{id}",
          "  postgres_message_count",
          "  recent_user_messages (last 3)",
          "  active_topic  (Redis value)",
          "  used_chunk_ids_count (Redis set)",
          "",
          "GET /devtools/jobs",
          "  In-memory active job states"], "OVR")

    card(slide, 0.20 + cw1 + gap, cy, cw2, ch, PURPLE, FILL_PURPLE, "Feature Flags",
         ["GET  /devtools/settings",
          "PATCH /devtools/settings",
          "",
          "enable_query_rewrite",
          "  Expands with history",
          "",
          "enable_hybrid_retrieval",
          "  Dense + BM25 search",
          "",
          "force_detailed_retrieval",
          "  Bypass mode limits",
          "",
          "rag_retrieval_mode",
          "  auto / semantic",
          "  keyword / hybrid",
          "",
          "rag_collection_name",
          "  Switch PGVector target"], "FLAGS")


def s_dbs(slide):
    cy = chrome(slide,
                "Database Visibility Panel",
                "GET /devtools/dbs/{db_id}/tables  |  GET /devtools/dbs/{db_id}/records?table=&limit=&offset=")
    cw = (SW - 0.58) / 4; ch = SH - cy - 0.28; gap = 0.18

    cards_data = [
        (BLUE,   FILL_BLUE,   "PGV",  "rag_db  (pgvector)",
         ["PostgreSQL + pgvector",
          "",
          "Tables:",
          "  langchain_pg_embedding",
          "  langchain_pg_collection",
          "",
          "Records browser:",
          "  Paginated limit/offset",
          "  embedding column",
          "  auto-excluded (too large)",
          "",
          "Stores semantic vectors",
          "384-dim BGE-M3 embeddings"]),
        (PURPLE, FILL_PURPLE, "CHAT", "chat_db  (postgres)",
         ["PostgreSQL chat memory",
          "",
          "Tables:",
          "  chat_messages",
          "  chat_sessions",
          "  session_topic_hints",
          "  rag_job_runs",
          "  rag_audit_log",
          "  retrieval_feedback",
          "",
          "Stores conversation",
          "history, audit events",
          "and user feedback"]),
        (GREEN,  FILL_GREEN,  "RDS",  "redis  (key-value)",
         ["Redis key browser",
          "",
          "Types shown:",
          "  string -> raw value",
          "  list   -> list[N]",
          "  set    -> set[N]",
          "  zset   -> zset[N]",
          "  hash   -> hash[N]",
          "",
          "Namespaces:",
          "  rag:*   retrieval cache",
          "  abort:* stream kill flags",
          "  user:*  session state"]),
        (ORANGE, FILL_ORANGE, "S3",   "minio  (object store)",
         ["MinIO object browser",
          "",
          "Fields shown:",
          "  object  (full path/key)",
          "  size    (bytes)",
          "  last_modified (ISO)",
          "",
          "Default bucket:",
          "  chat-ui-documents",
          "",
          "Lists ALL objects",
          "recursively, paginated",
          "Returns total count"]),
    ]
    for i, (ac, fi, ico, ttl, buls) in enumerate(cards_data):
        card(slide, 0.20 + i * (cw + gap), cy, cw, ch, ac, fi, ttl, buls, ico)


def s_runtime(slide):
    cy = chrome(slide,
                "Runtime Status  -  System Health Monitoring",
                "GET /devtools/runtime  |  GET /devtools/models/active  |  Per-mode readiness check")
    gw = (SW - 0.58) / 2; gh = (SH - cy - 0.36) / 2; gap = 0.18

    data = [
        (BLUE,   FILL_BLUE,   "GPU",  "GPU & Compute",
         ["Device name (NVIDIA RTX...)",
          "VRAM total / free / used (GB)",
          "CUDA version from torch build",
          "",
          "Determines mode selection:",
          "  base -> GPU if VRAM fits",
          "  lite -> CPU always",
          "",
          "Zero VRAM: CPU-only mode"]),
        (RED,    FILL_RED,    "MQ",   "RabbitMQ Broker",
         ["Broker reachability ping",
          "",
          "Queue depths:",
          "  default      general tasks",
          "  rag_ingest   PDF pipeline",
          "  rag_high_mem unstructured",
          "",
          "High depth = upload backlog"]),
        (GREEN,  FILL_GREEN,  "WRK",  "Celery Workers",
         ["Active worker list",
          "Task count per instance",
          "Queue assignments",
          "",
          "Worker types:",
          "  default      light async",
          "  rag_ingest   PDF chunking",
          "  rag_high_mem hi_res OCR",
          "",
          "Zero workers = stalled pipeline"]),
        (PURPLE, FILL_PURPLE, "MDL",  "Active Model Check",
         ["Per mode: base / lite / net",
          "  configured_model_id",
          "  effective_model_id",
          "  type: gguf / hf / net",
          "  path: local file/cache",
          "  ready: file exists?",
          "  loaded: in RAM now?",
          "  error: reason if not ready",
          "",
          "net: provider + key present?"]),
    ]
    for i, (ac, fi, ico, ttl, buls) in enumerate(data):
        ci = i % 2; ri = i // 2
        card(slide, 0.20 + ci * (gw + gap),
             cy + ri * (gh + gap * 0.55),
             gw, gh, ac, fi, ttl, buls, ico)


def s_reset(slide):
    cy = chrome(slide,
                "System Reset Controls  -  Admin Destructive Operations",
                "Guarded by env flag + confirmation phrase  |  Each layer can be reset independently or all at once")

    # Warning box
    wh = 0.88
    w = rect(slide, 0.20, cy, SW - 0.40, wh, fill=FILL_RED, border=RED, bpt=2.0)
    wtf = w.text_frame; wtf.word_wrap = True
    wtf.margin_left = Pt(12); wtf.margin_top = Pt(8)
    para(wtf, "SAFETY REQUIREMENTS  -  Both must be true before any reset fires:", size=12, color=RED, bold=True, spc_after=3)
    para(wtf, "  1.  Env var  CHAT_UI_ENABLE_DESTRUCTIVE_DEVTOOLS=1  must be set on the server.", size=11, color=BODY)
    para(wtf, "  2.  Request body must contain:   confirm = 'DELETE_EVERYTHING'", size=11, color=WHITE, bold=True)

    # Reset step flow
    fy = cy + wh + 0.20
    fh = SH - fy - 0.28
    steps = [
        {"title": "/reset/rag",
         "body":  ["TRUNCATE", "pg_embedding", "pg_collection", "clears all vectors"]},
        {"title": "/reset/chat",
         "body":  ["TRUNCATE", "chat_messages", "chat_sessions", "topics + docs"]},
        {"title": "/reset/redis",
         "body":  ["Delete rag:*", "and abort:* keys", "or full flushdb", "wipe_all=True"]},
        {"title": "/reset/minio",
         "body":  ["Remove ALL objects", "in target bucket", "minio_bucket", "field required"]},
        {"title": "/reset/all",
         "body":  ["Chains all 4", "resets in order", "Returns combined", "result per target"]},
    ]
    flow(slide, fy, fh, steps, [BLUE, PURPLE, GREEN, ORANGE, RED], label_size=12)


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

SLIDE_BUILDERS = [
    ("Developer Mode Overview",           s_overview),
    ("LLM Model Modes",                   s_model_modes),
    ("Model Download & Registry API",     s_model_api),
    ("RAG Debug Pipeline",                s_rag_debug),
    ("MVP Analytics Command Center",      s_mvp),
    ("Upload Job Inspector",              s_upload),
    ("User Management & Provisioning",    s_users),
    ("RAG Overrides & Session State",     s_overrides),
    ("Database Visibility Panel",         s_dbs),
    ("Runtime Status",                    s_runtime),
    ("System Reset Controls",             s_reset),
]


def main():
    global _KAVIN_LOGO_BYTES

    print(f"Opening: {PPT_PATH}")
    prs = Presentation(str(PPT_PATH))
    print(f"  Existing slides: {len(prs.slides)}")

    # Extract KAVIN logo from slide 1 (last picture shape found)
    try:
        sl1 = list(prs.slides)[0]
        for shape in sl1.shapes:
            if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE = 13
                _KAVIN_LOGO_BYTES = shape.image.blob
        if _KAVIN_LOGO_BYTES:
            print(f"  Extracted KAVIN logo: {len(_KAVIN_LOGO_BYTES)} bytes")
        else:
            print("  KAVIN logo not found, using text badge")
    except Exception as e:
        print(f"  Logo extract failed: {e}")

    print()
    for name, builder in SLIDE_BUILDERS:
        print(f"  Building: {name}")
        sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank
        builder(sl)

    total = len(prs.slides)
    print(f"\n  Total slides: {total}")

    try:
        prs.save(str(PPT_PATH))
        print(f"  Saved -> {PPT_PATH}")
    except PermissionError:
        alt = PPT_PATH.parent / "KavinBase_Final.pptx"
        prs.save(str(alt))
        print(f"  (File locked) Saved -> {alt}")


if __name__ == "__main__":
    main()
