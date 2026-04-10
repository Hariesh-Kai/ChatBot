import base64
import zlib
import requests
import io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PPT = BASE_DIR / "KavinBase_Presentation_Updated.pptx"


def asset_path(*relative_candidates):
    for candidate in relative_candidates:
        if not candidate:
            continue
        path = BASE_DIR / candidate
        if path.exists():
            return path
    return None

def get_mermaid_image(mermaid_text):
    try:
        compressed = zlib.compress(mermaid_text.encode('utf-8'), 9)
        b64 = base64.urlsafe_b64encode(compressed).decode('ascii')
        url = f"https://kroki.io/mermaid/png/{b64}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return io.BytesIO(response.content)
        else:
            print(f"Failed to fetch image: {response.text}")
            return None
    except Exception as e:
        print(f"Error fetching mermaid image: {e}")
        return None

prs = Presentation()

def add_slide(prs, title_text, content_bullets=None, image_stream=None, image_path=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    
    if (image_stream or image_path) and not content_bullets:
        slide_layout = prs.slide_layouts[5] # Title Only
        
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    if content_bullets:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_bullets[0]
        tf.paragraphs[0].font.size = Pt(16)
        for bullet in content_bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.level = 0
            
    if image_stream or image_path:
        left = Inches(1)
        top = Inches(2.2)
        height = Inches(4.5)
        if content_bullets:
             top = Inches(4.0)
             height = Inches(3.0)
             left = Inches(1.5)
        try:
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), left, top, height=height)
            elif image_stream:
                slide.shapes.add_picture(image_stream, left, top, height=height)
        except Exception as e:
            print(f"Error adding picture: {e}")
        
    return slide

# 1. Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "KavinBase System Overview"
subtitle.text = "End-to-End Artificial Intelligence Document Pipeline\nArchitecture & Implementation"

# 2. SYSTEM ARCHITECTURE
arch_image = asset_path("System_Architecture_Graph.png", "tmp/System_Architecture_Graph.png")
add_slide(prs, "System Architecture", ["The system acts as a hybrid synchronous-asynchronous platform:", "FastAPI backend handles high-speed REST interactions and SSE streaming limits.", "Celery offloads CPU-intensive PDF chunking and bounding-box extractions.", "Data propagates across multiple purpose-built persistence layers (MinIO, PostgreSQL PGVector, Redis)."], image_path=arch_image)

# 3. INGESTION PIPELINE
ingest_image = asset_path("Ingestion_Pipeline_Graph.png", "tmp/Ingestion_Pipeline_Graph.png")
add_slide(prs, "Ingestion Pipeline", ["Documents undergo a rigorous processing workflow:", "1. Uploads securely pushed to an S3-compatible MinIO datastore.", "2. Heavy lifting is assigned to asynchronous Celery tasks to prevent UI blockage.", "3. Unstructured parses complex layouts, while embeddings transform chunks into multi-dimensional vectors for semantic search."], image_path=ingest_image)

# 4. QUERY / CHAT PIPELINE
query_image = asset_path("Query_Pipeline_Graph.png", "tmp/Query_Pipeline_Graph.png")
add_slide(prs, "Query Pipeline", ["Hybrid intelligence retrieval system:", "Queries execute concurrently: performing semantic search alongside standard relational keyword hits.", "FlashRank takes the overlapping chunk results, applies Cross-Encoding models, and reranks them to ensure optimum context fidelity.", "Local LLMs then stream generating responses chunk by chunk over HTTP connections back to the Next.js UI."], image_path=query_image)

# 4B. CONCEPTUAL MODEL WORKFLOW IMAGE
model_img = asset_path("model_workflow_conceptual.png", "tmp/model_workflow_conceptual.png")
add_slide(prs, "AI Intelligence Workflow", ["A high-level abstraction of the RAG (Retrieval-Augmented Generation) pipeline.", "Demonstrates semantic data streaming off vectors and interacting with the neural LLM endpoint."], image_path=model_img)

# 5. FRONTEND: Auth & Next.js Architecture
auth_img = asset_path("dashboard_real.png", "tmp/dashboard_real.png")
add_slide(prs, "Frontend: Application Dashboard", ["The Dashboard serves as the main command center.", "Pulls real-time project metrics directly from the PostgreSQL instance.", "Secure JWT verification protects layout access across routes."], image_path=auth_img)

# 6. FRONTEND: Chat Interface
chat_img = asset_path("chat_real.png", "tmp/chat_real.png")
add_slide(prs, "Frontend: Chat Workspace", ["Real-time intelligence retrieval center.", "Maintains continuous SSE (Server-Sent Events) links with the LLM backend for token-by-token rendering.", "Embeds technical Markdown, Tables, and citations safely."], image_path=chat_img)

# 7. FRONTEND: Projects & Upload
proj_img = asset_path("projects_real.png", "tmp/projects_real.png")
add_slide(prs, "Frontend: Projects Hub", ["Document indexing and status tracking interface.", "Monitors active celery worker tasks and asynchronous indexing via simple UI indicators."], image_path=proj_img)

# 8. BACKEND: Swagger & API
api_img = asset_path("api_real.png", "tmp/api_real.png")
add_slide(prs, "Backend: API Exposure", ["The FastAPI backend strictly types payloads using Pydantic.", "Provides a fully interactive Swagger documentation interface out of the box.", "Modular routes separate Chat, Authentication, and Upload workflows seamlessly."], image_path=api_img)

# 9. BACKEND STORAGE SUMMARY
add_slide(prs, "Storage Ecosystem", [
    "MinIO (Object Stores): Retains raw PDFs and binary blobs isolated from the database.",
    "PostgreSQL + PGVector (Relational / Vector): Handles user schemas, project models, and the ~384 dimensional vector indexes.",
    "Redis (Key-Value Keyring): Serves as the high-speed Cache layer tracking streaming states and broker messaging."
])

# Save File
try:
    prs.save(str(OUTPUT_PPT))
    print(f"Successfully generated {OUTPUT_PPT}")
except Exception as e:
    print(f"Failed to save PPT: {e}")
