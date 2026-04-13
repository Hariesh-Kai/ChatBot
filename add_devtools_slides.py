"""
add_devtools_slides.py
Appends Developer Mode feature slides to the existing KavinBase_Presentation_Updated.pptx.
Existing slides are NOT changed — only new slides are appended.
"""

import base64
import zlib
import requests
import io
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR = Path(__file__).resolve().parent
INPUT_PPT  = BASE_DIR / "KavinBase_Presentation_Updated.pptx"
OUTPUT_PPT = BASE_DIR / "KavinBase_Presentation_DevTools.pptx"  # new file — original untouched

# ──────────────────────────────────────────────
# Mermaid → PNG via kroki.io
# ──────────────────────────────────────────────
def get_mermaid_image(mermaid_text: str):
    try:
        compressed = zlib.compress(mermaid_text.encode("utf-8"), 9)
        b64 = base64.urlsafe_b64encode(compressed).decode("ascii")
        url = f"https://kroki.io/mermaid/png/{b64}"
        resp = requests.get(url, timeout=20)
        if resp.status_code == 200:
            return io.BytesIO(resp.content)
        print(f"  [kroki] HTTP {resp.status_code}: {resp.text[:200]}")
        return None
    except Exception as e:
        print(f"  [kroki] Error: {e}")
        return None


# ──────────────────────────────────────────────
# Slide helpers
# ──────────────────────────────────────────────
ACCENT = RGBColor(0x1A, 0x73, 0xE8)   # Google-blue accent

def add_slide(prs, title_text, bullets=None, image_stream=None, image_path=None):
    """Add a slide with optional bullet list and/or image."""
    if (image_stream or image_path) and not bullets:
        layout = prs.slide_layouts[5]   # Title Only
    else:
        layout = prs.slide_layouts[1]   # Title and Content

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text

    if bullets:
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = bullets[0]
        tf.paragraphs[0].font.size = Pt(16)
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(16)
            p.level = 0

    if image_stream or image_path:
        left   = Inches(1)
        top    = Inches(2.2)
        height = Inches(4.5)
        if bullets:
            top    = Inches(4.0)
            height = Inches(3.0)
            left   = Inches(1.5)
        try:
            if image_path and Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), left, top, height=height)
            elif image_stream:
                slide.shapes.add_picture(image_stream, left, top, height=height)
        except Exception as e:
            print(f"  [image] Error adding picture: {e}")

    return slide


def add_split_slide(prs, title_text, left_bullets, right_bullets):
    """
    Two-column text slide — left column bullets | right column bullets.
    """
    layout = prs.slide_layouts[5]   # Title Only (we draw text boxes manually)
    slide  = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    col_top    = Inches(1.6)
    col_height = Inches(5.2)
    col_w      = Inches(4.4)
    gap        = Inches(0.2)
    left_x     = Inches(0.3)
    right_x    = left_x + col_w + gap

    def _add_text_box(x, y, w, h, lines):
        txb = slide.shapes.add_textbox(x, y, w, h)
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = Pt(15)
            run.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
            if line.startswith("●") or line.startswith("•"):
                p.level = 1

    _add_text_box(left_x,  col_top, col_w, col_height, left_bullets)
    _add_text_box(right_x, col_top, col_w, col_height, right_bullets)
    return slide


# ──────────────────────────────────────────────
# Mermaid diagrams
# ──────────────────────────────────────────────

MERMAID_MODEL_SELECTION = """
flowchart TD
    Q([User Query]) --> MS[Model Selector]
    MS --> B{Mode?}
    B -->|base| HF[HuggingFace\\nTransformer\\ne.g. Qwen2.5-7B]
    B -->|lite| GGUF[GGUF Quantized\\nllama.cpp\\ne.g. Qwen2.5-3B Q4]
    B -->|net| NET[Cloud API\\ne.g. OpenAI / Gemini]
    HF --> OUT([LLM Response])
    GGUF --> OUT
    NET --> OUT
    style Q fill:#1A73E8,color:#fff
    style OUT fill:#34A853,color:#fff
    style HF fill:#fbbc04,color:#000
    style GGUF fill:#ea4335,color:#fff
    style NET fill:#4285f4,color:#fff
"""

MERMAID_RAG_DEBUG = """
flowchart LR
    Q([Dev Question]) --> N[normalize_text]
    N --> IC{Intent\\nClassifier}
    IC -->|fact_lookup\\ncomparison\\netc.| RW[Query Rewriter\\noptional]
    RW --> KW[Keyword\\nExtractor]
    KW --> RT[retrieve_rag_context\\nPGVector + Hybrid]
    RT --> CH([Top N Chunks\\nreturned])
    style Q fill:#1A73E8,color:#fff
    style CH fill:#34A853,color:#fff
    style IC fill:#fbbc04,color:#000
    style RT fill:#ea4335,color:#fff
"""

MERMAID_USER_PROVISION = """
flowchart TD
    A([Admin POST /devtools/users]) --> CU[create_user\\npassword hash + role]
    CU --> PU[_provision_user_resources]
    PU --> PG[Ensure PostgreSQL DB\\nchat_ui_username]
    PU --> MB[Ensure MinIO Bucket\\nchat-ui-username]
    PU --> RD[Redis Namespace\\nuser:username:]
    PG & MB & RD --> SR[set_user_resources\\npersist metadata]
    SR --> OK([User Ready])
    style A fill:#1A73E8,color:#fff
    style OK fill:#34A853,color:#fff
    style PG fill:#4285f4,color:#fff
    style MB fill:#ea4335,color:#fff
    style RD fill:#fbbc04,color:#000
"""

MERMAID_MVP_OVERVIEW = """
flowchart LR
    OV([GET /devtools/mvp/overview]) --> RT[Runtime Status\\nGPU · RabbitMQ · Workers]
    OV --> IN[Ingestion Funnel\\nWaiting · Processing · Ready · Error]
    OV --> RQ[RAG Quality\\nHigh · Medium · Low · Cache · Latency]
    OV --> FB[Feedback Analytics\\nPositive · Negative · Label Distribution]
    style OV fill:#1A73E8,color:#fff
    style RT fill:#ea4335,color:#fff
    style IN fill:#fbbc04,color:#000
    style RQ fill:#34A853,color:#fff
    style FB fill:#4285f4,color:#fff
"""

# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────

def main():
    print(f"Opening: {INPUT_PPT}")
    prs = Presentation(str(INPUT_PPT))
    print(f"  Existing slides: {len(prs.slides)}")

    # ── SLIDE A: Developer Mode Introduction ─────────────────────────────────
    print("Adding slide: Developer Mode — Overview")
    add_slide(prs, "Developer Mode — Overview", [
        "Kavin AI ships a built-in /devtools admin panel for deep system introspection.",
        "● All endpoints are admin-gated (JWT role check via require_admin).",
        "● Feature flags can be toggled live without restarting the server.",
        "● Covers: Model mgmt · RAG debug · User admin · DB browser · Reset controls.",
        "● Destructive operations (reset) require an env flag + confirmation phrase.",
        "● Router prefix: /devtools  |  Tags: Developer Tools",
    ])

    # ── SLIDE B: LLM Model Management — Three Modes ──────────────────────────
    print("Adding slide: LLM Model Management")
    print("  Fetching model-selection flowchart …")
    model_img = get_mermaid_image(MERMAID_MODEL_SELECTION)
    add_slide(prs, "LLM Model Management — Three Modes", [
        "base  →  HuggingFace Transformers (GPU / CPU).  e.g. Qwen2.5-7B-Instruct",
        "lite   →  GGUF Quantized via llama.cpp (CPU-friendly).  e.g. Qwen2.5-3B Q4",
        "net   →  Cloud API provider (OpenAI, Gemini).  Key stored in secrets store.",
        "Admin can download, register, test and delete models at runtime via REST.",
    ], image_stream=model_img)

    # ── SLIDE C: Model Download & Registry ───────────────────────────────────
    print("Adding slide: Model Download & Registry")
    add_split_slide(prs,
        "Model Download & Registry API",
        left_bullets=[
            "DOWNLOAD ENDPOINTS",
            "● POST /devtools/models/download",
            "  Auto-detects GGUF vs HF snapshot",
            "● POST /devtools/models/gguf/download",
            "  Direct URL → models/gguf/ dir",
            "● POST /devtools/models/hf/install",
            "  snapshot_download to HF cache",
            "",
            "REGISTER / TEST",
            "● POST /devtools/models/gguf/register",
            "  Register local .gguf path",
            "● POST /devtools/models/test",
            "  Run a quick generation sample",
        ],
        right_bullets=[
            "KNOWN SUPPORTED MODELS",
            "● Qwen/Qwen2.5-3B-Instruct-GGUF",
            "● Qwen/Qwen2.5-1.5B-Instruct-GGUF",
            "● Qwen/Qwen2.5-7B-Instruct-GGUF",
            "● AI-Engine/Llama-3.1-8B-GGUF",
            "● Qwen/Qwen2.5-3B-Instruct (HF)",
            "● Qwen/Qwen2.5-7B-Instruct (HF)",
            "",
            "AUX / EMBEDDING MODELS",
            "● BAAI/bge-m3  → embedding",
            "● facebook/bart-large-mnli → intent",
            "● microsoft/table-transformer",
            "● unstructuredio/yolo_x_layout",
        ]
    )

    # ── SLIDE D: RAG Debug & Retrieval Testing ───────────────────────────────
    print("Adding slide: RAG Debug & Retrieval Testing")
    print("  Fetching RAG debug flowchart …")
    rag_img = get_mermaid_image(MERMAID_RAG_DEBUG)
    add_slide(prs, "RAG Debug & Retrieval Testing", [
        "POST /devtools/intent   →  Classify query intent (fact_lookup, comparison, …)",
        "POST /devtools/rewrite  →  Show how a question is rewritten with history.",
        "POST /devtools/keywords →  Inspect keywords extracted for SQL hybrid search.",
        "POST /devtools/retrieve →  Full RAG retrieval dry-run, returns top N chunks.",
    ], image_stream=rag_img)

    # ── SLIDE E: MVP Analytics Command Center ────────────────────────────────
    print("Adding slide: MVP Analytics Command Center")
    print("  Fetching MVP overview flowchart …")
    mvp_img = get_mermaid_image(MERMAID_MVP_OVERVIEW)
    add_slide(prs, "MVP Analytics Command Center", [
        "GET /devtools/mvp/overview  (window_hours: 1–168h)",
        "● Ingestion Funnel: waiting → processing → ready | error + success rate.",
        "● RAG Quality: high / medium / low distribution + avg latency + cache hit %.",
        "● Feedback: positive / negative counts, avg score, label distribution.",
        "● Runtime: GPU VRAM, RabbitMQ broker health, Celery worker queues.",
    ], image_stream=mvp_img)

    # ── SLIDE F: Recent Upload Jobs & Preprocessing Preview ──────────────────
    print("Adding slide: Upload Job Inspector")
    add_slide(prs, "Upload Job Inspector & Preprocessing Preview", [
        "GET /devtools/uploads/recent  →  Last N ingestion jobs with full metadata.",
        "● Fields: job_id, status, progress %, progress_label, error message.",
        "● Company document ID, revision, source file, RAG preprocessor used.",
        "● preview_available flag — links to extracted element JSON artifacts.",
        "● Segregation summary: raw_elements, filtered_elements, removed, chunks.",
        "● artifact_only_preview flag when PDF is gone but cached artifacts exist.",
    ])

    # ── SLIDE G: User Management & Resource Provisioning ─────────────────────
    print("Adding slide: User Management & Resource Provisioning")
    print("  Fetching user provisioning flowchart …")
    user_img = get_mermaid_image(MERMAID_USER_PROVISION)
    add_slide(prs, "User Management & Resource Provisioning", [
        "POST /devtools/users   →  Create user + auto-provision all storage resources.",
        "PATCH /devtools/users/disable  →  Enable or disable a user account.",
        "PATCH /devtools/users/password →  Admin password reset (no old pwd needed).",
        "PATCH /devtools/users/role     →  Promote / demote user to admin.",
        "DELETE /devtools/users →  Remove user record from auth store.",
    ], image_stream=user_img)

    # ── SLIDE H: RAG Overrides & Session Control ──────────────────────────────
    print("Adding slide: RAG Overrides & Session Control")
    add_split_slide(prs,
        "RAG Overrides & Session State Inspector",
        left_bullets=[
            "RAG OVERRIDE CONTROLS",
            "● GET /devtools/rag/overrides",
            "  List all active session/user overrides",
            "● POST /devtools/rag/disable",
            "  Disable RAG for session or user",
            "● POST /devtools/rag/enable",
            "  Re-enable RAG for session or user",
            "",
            "Use case: Force direct LLM answers",
            "without document context retrieval,",
            "per user or per chat session.",
        ],
        right_bullets=[
            "SESSION STATE INSPECTOR",
            "● GET /devtools/session-state/{id}",
            "  postgres_message_count",
            "  recent_user_messages (last 3)",
            "  active_topic  (Redis)",
            "  used_chunk_ids_count (Redis)",
            "",
            "FEATURE FLAGS (settings)",
            "● GET/PATCH /devtools/settings",
            "  enable_query_rewrite",
            "  enable_hybrid_retrieval",
            "  force_detailed_retrieval",
            "  rag_retrieval_mode",
        ]
    )

    # ── SLIDE I: Database Visibility Panel ───────────────────────────────────
    print("Adding slide: Database Visibility Panel")
    add_slide(prs, "Database Visibility Panel", [
        "GET /devtools/dbs                   →  List all registered data sources.",
        "GET /devtools/dbs/{db_id}/tables    →  Tables/keys/objects per source.",
        "GET /devtools/dbs/{db_id}/records   →  Paginated row/key/object browser.",
        "● rag_db    (PostgreSQL + pgvector)  — langchain_pg_embedding, collections.",
        "● chat_db   (PostgreSQL)             — chat_messages, sessions, topics.",
        "● redis     (Key-Value)              — token browser, type + value preview.",
        "● minio     (Object Store)           — bucket object list with size & date.",
        "Embedding vectors are automatically excluded from responses (too large).",
    ])

    # ── SLIDE J: Runtime Status ───────────────────────────────────────────────
    print("Adding slide: Runtime Status")
    add_slide(prs, "Runtime Status — System Health", [
        "GET /devtools/runtime  →  Full runtime snapshot for ops monitoring.",
        "● GPU: device name, total/free/used VRAM, CUDA version.",
        "● RabbitMQ: broker reachability, queue depths per worker type.",
        "● Workers: active Celery workers, task counts, queue assignments.",
        "● Software: Python version, torch version, llama-cpp version.",
        "GET /devtools/models/active  →  Per-mode model readiness check.",
        "  Reports: type (gguf/hf/net), path, loaded-into-memory flag, errors.",
    ])

    # ── SLIDE K: System Reset Controls ───────────────────────────────────────
    print("Adding slide: System Reset Controls")
    add_slide(prs, "System Reset Controls (Admin-Gated)", [
        "All reset endpoints require:  env CHAT_UI_ENABLE_DESTRUCTIVE_DEVTOOLS=1",
        "Plus confirmation phrase in body:  confirm = 'DELETE_EVERYTHING'",
        "─────────────────────────────────────────────────────",
        "POST /devtools/reset/rag    →  TRUNCATE langchain_pg_embedding + collections.",
        "POST /devtools/reset/chat   →  TRUNCATE chat_messages, sessions, topics.",
        "POST /devtools/reset/redis  →  Delete rag:* and abort:* keys (or full flush).",
        "POST /devtools/reset/minio  →  Remove all objects in the target bucket.",
        "POST /devtools/reset/all    →  Chains all four resets in sequence.",
    ])

    # ── Save ──────────────────────────────────────────────────────────────────
    try:
        prs.save(str(OUTPUT_PPT))
        total = len(prs.slides)
        print(f"\nSaved -> {OUTPUT_PPT}")
        print(f"Total slides now: {total}")
    except Exception as e:
        print(f"\nFailed to save: {e}")


if __name__ == "__main__":
    main()
