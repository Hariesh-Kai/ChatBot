"""
analyze_ppt.py — Reads every shape from every slide and prints
background colors, fill colors, text runs, fonts, sizes, positions.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

PPT = Path(r"D:\KavinBase\KavinBase_Presentation_Updated.pptx")
prs = Presentation(str(PPT))

SW = prs.slide_width.inches
SH = prs.slide_height.inches
print(f"Slide size: {SW:.3f}\" x {SH:.3f}\"")
print(f"Total slides: {len(prs.slides)}\n")

def rgb_str(color):
    try:
        r, g, b = color.rgb
        return f"#{r:02X}{g:02X}{b:02X}"
    except Exception:
        return str(color)

def emu_to_in(v):
    return round(v / 914400, 3)

def dump_fill(fill, label="fill"):
    try:
        ft = fill.type
        if ft is None:
            return f"{label}: none"
        if str(ft) == "SOLID (1)":
            return f"{label}: solid {rgb_str(fill.fore_color)}"
        return f"{label}: {ft}"
    except Exception as e:
        return f"{label}: err({e})"

def dump_line(line, label="line"):
    try:
        w = round(line.width / 12700, 1) if line.width else 0
        col = rgb_str(line.color) if line.color and line.color.type else "none"
        return f"{label}: color={col} width={w}pt"
    except Exception as e:
        return f"{label}: err({e})"

# Only inspect slides 1-5 deeply + print slide backgrounds
for si, slide in enumerate(list(prs.slides)[:8], 1):
    print(f"{'='*60}")
    print(f"SLIDE {si}")

    # Background
    bg = slide.background
    try:
        bf = bg.fill
        print(f"  Background: {dump_fill(bf, 'bg')}")
    except Exception as e:
        print(f"  Background: err({e})")

    for shape in slide.shapes:
        stype = shape.shape_type
        x = emu_to_in(shape.left or 0)
        y = emu_to_in(shape.top or 0)
        w = emu_to_in(shape.width or 0)
        h = emu_to_in(shape.height or 0)
        name = shape.name

        fill_info = ""
        line_info = ""
        try:
            fill_info = dump_fill(shape.fill)
        except Exception:
            pass
        try:
            line_info = dump_line(shape.line)
        except Exception:
            pass

        texts = []
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text.strip()
                    if not txt:
                        continue
                    fc = "?"
                    sz = "?"
                    bold = False
                    fname = "?"
                    try:
                        fc = rgb_str(run.font.color)
                        sz = round(run.font.size / 12700, 1) if run.font.size else "inherit"
                        bold = run.font.bold
                        fname = run.font.name or "inherit"
                    except Exception:
                        pass
                    texts.append(f"'{txt[:40]}' font={fname} sz={sz}pt color={fc} bold={bold}")

        print(f"  [{name}] pos=({x}\",{y}\") size={w}\"x{h}\"")
        if fill_info:
            print(f"    {fill_info}")
        if line_info and "none" not in line_info and "err" not in line_info:
            print(f"    {line_info}")
        for t in texts[:4]:
            print(f"    TEXT: {t}")

print("\nDone.")
